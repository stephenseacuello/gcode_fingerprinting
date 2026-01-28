#!/usr/bin/env python3
"""
Create G-code Fingerprinting v22 Results PowerPoint Presentation - V4

V4 changes from v3:
  1. Feature breakdown corrected: 16 sensors x 17 modalities (272) + 24 machine features = 296
  2. Slide #18 (sensor rankings): full 12-sensor table restored to core deck
  3. Slides #22-23 (modality): all 17 modalities listed explicitly (Ax, Ay, Az, etc.)
  4. Writing quality improved: technical, precise, concise language throughout
  5. Vocabulary corrected: 2,498 tokens (not "2,000+")
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Constants
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 112, 192)
GREEN = RGBColor(0, 128, 0)
RED = RGBColor(192, 0, 0)
ORANGE = RGBColor(255, 153, 0)
GRAY = RGBColor(89, 89, 89)
GOLD = RGBColor(255, 215, 0)

# Global slide counter for slide numbers
slide_counter = [0]

# Complete sensor ablation results (v21 pipeline)
SENSOR_RESULTS = [
    ("frame_r2", 93.50, 77.07, 92.98),
    ("xa_motor", 93.35, 76.59, 93.78),
    ("frame_r1", 93.04, 74.63, 93.74),
    ("y_bed__1", 92.93, 74.80, 93.47),
    ("spindle2", 92.54, 73.01, 92.98),
    ("frame_l2", 92.35, 72.03, 91.32),
    ("frame_l3", 92.28, 71.22, 91.79),
    ("frame_b2", 92.20, 71.22, 93.28),
    ("y_bed__2", 92.12, 71.54, 94.40),
    ("y_bed__3", 91.55, 72.52, 91.90),
    ("y_bed__4", 91.09, 69.92, 93.17),
    ("spindle1", 88.07, 58.54, 91.17),
]

# Complete modality ablation results (v21 pipeline, baseline 80.69%)
MODALITY_ONLY_RESULTS = [
    ("RMS (Computed)", 90.44, 67.48, 40),
    ("Gyroscope", 89.06, 61.30, 72),
    ("Magnetometer", 87.92, 58.21, 72),
    ("Color", 86.69, 55.61, 88),
    ("Accelerometer", 86.50, 51.71, 72),
    ("Environmental", 82.72, 39.67, 72),
]

MODALITY_WITHOUT_RESULTS = [
    ("Without Environmental", 84.36, 45.85, 248),
    ("Without Computed (RMS)", 81.91, 37.40, 280),
    ("Without Accelerometer", 80.96, 37.24, 248),
    ("Without Gyroscope", 81.57, 39.02, 248),
    ("Without Magnetometer", 81.45, 36.42, 248),
    ("Without Color", 79.01, 37.56, 232),
]

# V22 ensemble results
ENSEMBLE_RESULTS = [
    (42, 84.28, 50.57, 85.83),
    (123, 82.98, 50.08, 84.76),
    (456, 83.75, 47.80, 85.37),
]

# Per-modality high correlations
HIGH_CORRELATIONS = [
    ("Pressure", "frame_r1 - frame_b2", "1.000"),
    ("Temperature", "frame_r1 - frame_b2", "1.000"),
    ("Ax", "frame_r1 - frame_b2", "1.000"),
    ("Pressure", "frame_r2 - frame_l2", "1.000"),
    ("RMS", "frame_r2 - y_bed__3", "0.997"),
    ("RMS", "frame_r1 - frame_b2", "0.996"),
    ("My", "frame_r1 - frame_b2", "-0.990"),
    ("Mz", "frame_r1 - frame_b2", "0.988"),
    ("RMS", "frame_r2 - spindle2", "0.986"),
    ("Mx", "frame_r1 - frame_b2", "-0.980"),
    ("Gz", "frame_r1 - frame_b2", "0.973"),
]


def set_cell_text(cell, text, font_size=12, bold=False, color=None, align_center=True):
    """Set text in a table cell with formatting."""
    cell.text = str(text)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        if color:
            paragraph.font.color.rgb = color
        if align_center:
            paragraph.alignment = PP_ALIGN.CENTER


def add_slide_number(slide):
    """Add slide number to bottom-right corner."""
    slide_counter[0] += 1
    num_box = slide.shapes.add_textbox(
        Inches(12.0), Inches(7.05), Inches(1.0), Inches(0.35)
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(slide_counter[0])
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2), SLIDE_WIDTH, Inches(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.33), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

    add_slide_number(slide)
    return slide


def add_section_slide(prs, section_title, section_number=""):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    if section_number:
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_number
        p.font.size = Pt(80)
        p.font.bold = True
        p.font.color.rgb = RGBColor(100, 150, 200)
        p.alignment = PP_ALIGN.CENTER

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.33), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    add_slide_number(slide)
    return slide


def add_content_slide(prs, title, bullet_points=None):
    """Add a content slide with title and optional bullet points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    if bullet_points:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.33), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            if point.startswith("**") and point.endswith("**"):
                p.text = point.strip("*")
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = DARK_BLUE
                p.space_before = Pt(14)
            elif point.startswith("  - "):
                p.text = point.strip()
                p.font.size = Pt(16)
                p.level = 1
            elif point == "":
                p.text = ""
                p.space_before = Pt(6)
            else:
                p.text = f"• {point}"
                p.font.size = Pt(18)
                p.space_before = Pt(4)

    add_slide_number(slide)
    return slide


def add_table_slide(prs, title, headers, rows, highlight_rows=None, col_widths=None, font_size=11, subtitle=None):
    """Add a slide with a table."""
    slide = add_content_slide(prs, title)

    start_top = Inches(1.3)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.4))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = GRAY
        start_top = Inches(1.6)

    num_cols = len(headers)
    num_rows = len(rows) + 1

    left = Inches(0.3)
    width = Inches(12.7)
    row_height = 0.35
    height = Inches(row_height * min(num_rows, 16))

    table = slide.shapes.add_table(num_rows, num_cols, left, start_top, width, height).table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    else:
        col_width = 12.7 / num_cols
        for i in range(num_cols):
            table.columns[i].width = Inches(col_width)

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        set_cell_text(cell, header, font_size=font_size + 1, bold=True)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = RGBColor(255, 255, 255)

    highlight_rows = highlight_rows or []
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            is_highlight = row_idx in highlight_rows
            set_cell_text(cell, cell_data, font_size=font_size, bold=is_highlight)
            if is_highlight:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 242, 204)
            elif row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)

    return slide


def add_image_slide(prs, title, image_path, caption="", image_width=None):
    """Add a slide with an image."""
    slide = add_content_slide(prs, title)

    if os.path.exists(image_path):
        img_width = Inches(image_width) if image_width else Inches(9)
        img_left = (SLIDE_WIDTH - img_width) / 2
        img_top = Inches(1.4)
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)

    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_two_image_slide(prs, title, image1_path, image2_path, caption1="", caption2=""):
    """Add a slide with two images side by side."""
    slide = add_content_slide(prs, title)

    img_width = Inches(5.8)
    img_top = Inches(1.5)

    if os.path.exists(image1_path):
        slide.shapes.add_picture(image1_path, Inches(0.5), img_top, width=img_width)

    if os.path.exists(image2_path):
        slide.shapes.add_picture(image2_path, Inches(7), img_top, width=img_width)

    if caption1:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(5.8), Inches(0.4))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption1
        p.font.size = Pt(11)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    if caption2:
        cap_box = slide.shapes.add_textbox(Inches(7), Inches(6.5), Inches(5.8), Inches(0.4))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption2
        p.font.size = Pt(11)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_two_column_slide(prs, title, left_content, right_content):
    """Add a slide with two columns."""
    slide = add_content_slide(prs, title)

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.9), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if point.startswith("**"):
            p.text = point.strip("*")
            p.font.bold = True
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_BLUE
            p.space_before = Pt(10)
        elif point == "":
            p.text = ""
        else:
            p.text = f"• {point}"
            p.font.size = Pt(14)

    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.4), Inches(5.9), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if point.startswith("**"):
            p.text = point.strip("*")
            p.font.bold = True
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_BLUE
            p.space_before = Pt(10)
        elif point == "":
            p.text = ""
        else:
            p.text = f"• {point}"
            p.font.size = Pt(14)

    return slide


def create_presentation():
    """Create the full v22 presentation - V4."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    base_path = "/Users/stepheneacuello/Projects/gcode_fingerprinting"

    # ================================================================
    # CORE SLIDES
    # ================================================================

    # ========== SLIDE: Title ==========
    add_title_slide(
        prs,
        "G-code Fingerprinting",
        "2K Vocabulary Experiment Results  |  Version 22  |  January 2026"
    )

    # ========== SLIDE: Problem Statement ==========
    add_content_slide(
        prs,
        "Problem Statement",
        [
            "**Objective**",
            "  - Reconstruct CNC G-code instruction sequences from passive sensor measurements",
            "  - Enable non-invasive machine fingerprinting without controller access",
            "",
            "**Motivation**",
            "  - Manufacturing security: detect unauthorized toolpath modifications",
            "  - Quality assurance: verify executed operations match the intended CAM design",
            "  - Process monitoring: continuous non-contact observation of CNC behavior",
            "",
            "**Approach**",
            "  - Instrument a Bantam Explorer CNC with 16 Arduino Nano 33 BLE Sense units (12 operational, 4 non-functional)",
            "  - Record 17 modalities per sensor (3-axis accel/gyro/mag, pressure, temp, proximity, RGBA, RMS)",
            "  - Train a two-stage encoder-decoder to predict G-code token sequences from sensor tensors",
            "",
            "**Technical Challenge**",
            "  - Vocabulary of 2,498 tokens with 4-digit numeric precision",
            "  - Model must predict exact coordinate values (e.g., X-12.3456), not just command types",
        ]
    )

    # ========== SLIDE: Executive Summary ==========
    add_table_slide(
        prs,
        "Executive Summary: Key Results",
        ["Configuration", "Token Acc.", "Seq. Acc.", "Features"],
        [
            ["Best single sensor (frame_r2)*", "93.50%", "77.1%", "41"],
            ["All 16 sensors (v22 corrected)", "84.28%", "50.6%", "296"],
            ["All 16 sensors (without color)", "75.37%", "28.0%", "208"],
            ["Performance gap (single vs multi)", "-9.22%", "-26.5%", "+255"],
        ],
        highlight_rows=[0, 3],
        font_size=14,
        subtitle="*Single-sensor results from v21 pipeline; v22 re-run pending. 296 features include 68 from 4 non-functional sensors."
    )

    # ========== SECTION: Dataset & Architecture ==========
    add_section_slide(prs, "Dataset & Architecture", "01")

    # ========== SLIDE: Dataset Overview ==========
    add_two_column_slide(
        prs,
        "Dataset Overview",
        [
            "**Data Composition**",
            "145 aligned CSV files from real CNC operations",
            "3,160+ training sequences (64 timesteps each)",
            "2,498-token vocabulary (4-digit precision)",
            "",
            "**Operations (145 files)**",
            "Face operations: 40 files (27.6%)",
            "Pocket operations: 40 files (27.6%)",
            "Adaptive operations: 40 files (27.6%)",
            "Damaged operations: 15 files (10.3%)",
            "150025 variants: 10 files (6.9%)",
        ],
        [
            "**16 Deployed Sensors (12 Operational)**",
            "Frame: frame_r1, frame_r2, frame_b2, frame_l2, frame_l3",
            "Spindle: spindle1, spindle2",
            "Y-bed: y_bed__1, y_bed__2, y_bed__3, y_bed__4",
            "Motor: xa_motor",
            "Non-functional (4): frame_b1, frame_l1, z_gant_1, z_gant_2",
            "",
            "**17 Modalities per Sensor**",
            "Accelerometer: Ax, Ay, Az",
            "Gyroscope: Gx, Gy, Gz",
            "Magnetometer: Mx, My, Mz",
            "Environmental: Pressure, Temperature, Proximity",
            "Color: ColorR, ColorG, ColorB, ColorA",
            "Computed: RMS",
        ]
    )

    # ========== SLIDE: Feature Breakdown ==========
    add_table_slide(
        prs,
        "Feature Composition (296 Total)",
        ["Category", "Source", "Count"],
        [
            ["12 operational sensors", "12 x 17 modalities (Ax..RMS)", "204"],
            ["4 non-functional sensors", "4 x 17 modalities (frame_b1, frame_l1, z_gant_1, z_gant_2)", "68"],
            ["Motor electrical", "Voltage + current for spindle, X, Y, Z motors", "8"],
            ["Machine positions", "Desired (X,Y,Z) + measured (X,Y,Z) + line number", "7"],
            ["Controller signals", "Status, units, distance mode, feed, plane, motion mode, etc.", "9"],
            ["Total", "204 + 68 + 24 machine features", "296"],
        ],
        highlight_rows=[0, 5],
        font_size=12,
        subtitle="16 sensors deployed; 4 ceased functioning during data collection. All 296 features remain in the input vector."
    )

    # ========== SLIDE: Preprocessing Pipeline ==========
    add_content_slide(
        prs,
        "Preprocessing Pipeline",
        [
            "**Step 1: Temporal Alignment**",
            "  - Interpolate sensor + motor data onto controller (SR) timestamps",
            "  - Linear interpolation via np.interp() for consistent time base",
            "",
            "**Step 2: Feature Construction (296 features)**",
            "  - 16 sensors x 17 modalities (272) + 24 machine features (8 motor + 7 position + 9 controller) = 296",
            "",
            "**Step 3: G-code Tokenization (2,498 tokens)**",
            "  - Hybrid tokenizer: literal command tokens (G0, G1, M3) + bucketed numeric value tokens",
            "  - Canonicalization: strip comments, normalize command codes, round to 4-digit precision",
            "",
            "**Step 4: Windowing**",
            "  - 64 timesteps per window, stride 16 (75% overlap)",
            "  - Each window produces a sensor tensor [64, 296] + G-code token sequence",
            "",
            "**Step 5: Feature Selection and Cleaning**",
            "  - Remove >50% missing, zero-variance, and |r|>0.95 correlated features",
            "  - Forward-fill NaN, clip outliers (IQR), log-transform skewed features",
        ]
    )

    # ========== SLIDE: Preprocessing Pipeline (cont.) ==========
    add_two_column_slide(
        prs,
        "Preprocessing Pipeline (cont.)",
        [
            "**Step 6: Normalization**",
            "v21: Global StandardScaler on all data (data leakage)",
            "v22: StandardScaler on training data only",
            "Per-feature (each of 296 columns independently)",
            "",
            "**Step 7: Train / Val / Test Split**",
            "70% train: 93 files (2,966 samples)",
            "15% val: 21 files (614 samples)",
            "15% test: 21 files (615 samples)",
            "File-level split (no sample leakage)",
        ],
        [
            "**Step 8: Data Augmentation (training only)**",
            "Sensor noise (Gaussian, sigma=0.01)",
            "Temporal shift (+/-1 timestep)",
            "Magnitude scaling (0.98-1.02x)",
            "Time warping (non-linear distortion)",
            "Feature dropout (p=0.05 per channel)",
            "Temporal cutout (mask 1-2 timesteps)",
            "Class balancing (3x oversample rare G-codes)",
            "Mixup (blend similar operations, alpha=0.2)",
        ]
    )

    # ========== SLIDE: Experimental Setup Photo ==========
    setup_photo = f"{base_path}/outputs/archive/2024-12_best_runs/sensor_multihead_v3/visualizations/paper/figures/setup.png"
    sensor_setup = f"{base_path}/outputs/archive/2024-12_best_runs/sensor_multihead_v3/visualizations/paper/figures/sensor_setup.png"
    add_two_image_slide(
        prs,
        "Experimental Setup and Sensor Hardware",
        setup_photo,
        sensor_setup,
        "Bantam Explorer Desktop CNC with mounted Arduino Nano 33 BLE Sense sensors",
        "Sensor pipeline: 9-axis IMU, Pressure, Temperature, Proximity, Color (RGBA)"
    )

    # ========== SLIDE: Architecture Diagram ==========
    arch_path = f"{base_path}/outputs/archive/2024-12_best_runs/sensor_multihead_v3/visualizations/paper/figures/simplified_architecture.png"
    add_image_slide(
        prs,
        "Two-Stage Architecture",
        arch_path,
        "Stage 1: MM-DTAE-LSTM encoder classifies operation type (100% accuracy). Stage 2: Transformer multi-head decoder predicts G-code token sequences.",
        image_width=11
    )

    # ========== SLIDE: Training Curves ==========
    train_curves = f"{base_path}/outputs/archive/2024-12_best_runs/sensor_multihead_v3/visualizations/paper/figures/training_curves_4panel.png"
    add_image_slide(
        prs,
        "Training Dynamics: Encoder and Decoder",
        train_curves,
        "Encoder converges to 100% accuracy by epoch 30. Decoder reaches ~90% token accuracy via three-phase curriculum learning (P1/P2/P3).",
        image_width=7
    )

    # ========== SLIDE: Encoder Performance ==========
    encoder_cm = f"{base_path}/outputs/production/encoder/confusion_matrix.png"
    add_image_slide(
        prs,
        "Stage 1: Encoder Operation Classification (100% Accuracy)",
        encoder_cm,
        "Perfect classification of all 9 CNC operation types (face, pocket, adaptive, damaged, etc.) from sensor input alone.",
        image_width=8
    )

    # ========== SLIDE: Decoder Confusion Matrices (FIX #11 - moved here) ==========
    type_cm = f"{base_path}/outputs/real_confusion_matrices/confusion_type.png"
    add_image_slide(
        prs,
        "Stage 2: Token Type Classification (4 Classes)",
        type_cm,
        "Numeric: 99%, Special: 100%, Parameter: 89%, Command: 80%",
        image_width=7
    )

    # ========== SECTION: V22 Improvements ==========
    add_section_slide(prs, "V22 Improvements", "02")

    # ========== SLIDE: Data Leakage Fix ==========
    add_table_slide(
        prs,
        "V22 Improvement #1: Data Leakage Fix",
        ["Metric", "v21 (With Leakage)", "v22 (Corrected)", "Change"],
        [
            ["Multi-sensor token acc", "80.69%", "84.28%", "+3.59%"],
            ["Multi-sensor seq acc", "35.4%", "50.6%", "+15.2%"],
            ["Single sensor (frame_r2)", "93.50%", "Not yet re-run", "Pending"],
            ["Performance gap", "12.81%", "9.22%", "Narrowed"],
        ],
        highlight_rows=[0, 1],
        font_size=14,
        subtitle="Issue: StandardScaler fitted on all data. Fix: Fitted only on training data."
    )

    # ========== SLIDE: Model Stability ==========
    ensemble_rows = [[str(s), f"{t:.2f}%", f"{sq:.2f}%", f"{v:.2f}%"]
                     for s, t, sq, v in ENSEMBLE_RESULTS]
    ensemble_rows.append(["Mean", "83.67%", "49.48%", "85.32%"])
    ensemble_rows.append(["Std Dev", "+/-0.66%", "+/-1.52%", "+/-0.54%"])

    add_table_slide(
        prs,
        "V22 Improvement #2: Model Stability Across Seeds",
        ["Seed", "Token Acc.", "Seq. Acc.", "Val. Token"],
        ensemble_rows,
        highlight_rows=[3, 4],
        font_size=14,
        subtitle="Low variance (+/-0.66%) demonstrates robust, well-tuned architecture."
    )

    # ========== SLIDE: Color Ablation ==========
    add_table_slide(
        prs,
        "V22 Improvement #3: Color Sensors Are Essential",
        ["Configuration", "Features", "Token Acc.", "Seq. Acc."],
        [
            ["With Color (RGBA)", "296", "84.28%", "50.57%"],
            ["Without Color", "208", "75.37%", "27.97%"],
            ["Difference", "-88", "-8.91%", "-22.60%"],
        ],
        highlight_rows=[2],
        font_size=14,
        subtitle="Despite R, G, B, A being 97%+ correlated internally, they provide essential cross-modal information."
    )

    # ========== SECTION: Sensor Analysis ==========
    add_section_slide(prs, "Sensor Analysis", "03")

    # ========== SLIDE: Full Sensor Rankings (all 12) ==========
    sensor_rows = [[str(i+1), name, f"{tok:.2f}%", f"{seq:.2f}%", f"{val:.2f}%"]
                   for i, (name, tok, seq, val) in enumerate(SENSOR_RESULTS)]
    sensor_rows.append(["-", "All 16 sensors (v22)", "84.28%", "50.6%", "85.83%"])

    add_table_slide(
        prs,
        "Single Sensor Performance (v21 pipeline, 12 operational sensors)",
        ["Rank", "Sensor", "Test Token", "Test Seq.", "Val Token"],
        sensor_rows,
        highlight_rows=[0, 12],
        font_size=10,
        subtitle="Each sensor uses 41 features (17 modalities + 24 machine features). All 12 operational sensors individually outperform the 296-feature model (which includes 68 non-functional sensor features)."
    )

    # ========== SLIDE: Per-Modality Correlation ==========
    per_mod_corr = f"{base_path}/outputs/jan23_followup/correlations/per_modality_correlation_summary.png"
    add_image_slide(
        prs,
        "Per-Modality Sensor Correlations (Prof. Sodhi's Insight)",
        per_mod_corr,
        "Per-modality analysis reveals frame sensors are highly correlated (r=1.00) for Pressure, Temperature, Ax.",
        image_width=11
    )

    # ========== SLIDE: Sensor Correlation Heatmap ==========
    sensor_corr = f"{base_path}/outputs/analysis/sensor_correlation_heatmap.png"
    add_image_slide(
        prs,
        "Sensor Correlation Matrix (Averaged Across Modalities)",
        sensor_corr,
        "Original analysis: max |r| = 0.56, suggesting sensors are not redundant. But this masked per-modality patterns.",
        image_width=8
    )

    # ========== SECTION: Modality Analysis ==========
    add_section_slide(prs, "Modality Analysis", "04")

    # ========== SLIDE: Single Modality Performance ==========
    mod_rows_detailed = [
        ["RMS (Computed)", "RMS (1 ch/sensor)", "90.44%", "67.48%", "40"],
        ["Gyroscope", "Gx, Gy, Gz (3 ch/sensor)", "89.06%", "61.30%", "72"],
        ["Magnetometer", "Mx, My, Mz (3 ch/sensor)", "87.92%", "58.21%", "72"],
        ["Color", "ColorR, ColorG, ColorB, ColorA (4 ch/sensor)", "86.69%", "55.61%", "88"],
        ["Accelerometer", "Ax, Ay, Az (3 ch/sensor)", "86.50%", "51.71%", "72"],
        ["Environmental", "Pressure, Temperature, Proximity (3 ch/sensor)", "82.72%", "39.67%", "72"],
        ["All Modalities (v21)", "All 17 channels/sensor", "80.69%", "35.4%", "296"],
    ]

    add_table_slide(
        prs,
        "Single Modality Performance (v21 pipeline)",
        ["Modality Group", "Channels", "Token Acc.", "Seq. Acc.", "Features"],
        mod_rows_detailed,
        highlight_rows=[0, 6],
        font_size=11,
        subtitle="Each row uses only that modality across all 12 operational sensors. RMS alone (90.44%) outperforms all 17 modalities combined (80.69%)."
    )

    # ========== SLIDE: Leave-One-Out Modality ==========
    without_rows_detailed = [
        ["Without Environmental", "Remove: Pressure, Temp, Proximity", "84.36%", "45.85%", "248"],
        ["Without Computed (RMS)", "Remove: RMS", "81.91%", "37.40%", "280"],
        ["Without Gyroscope", "Remove: Gx, Gy, Gz", "81.57%", "39.02%", "248"],
        ["Without Magnetometer", "Remove: Mx, My, Mz", "81.45%", "36.42%", "248"],
        ["Without Accelerometer", "Remove: Ax, Ay, Az", "80.96%", "37.24%", "248"],
        ["Without Color", "Remove: ColorR, G, B, A", "79.01%", "37.56%", "232"],
    ]

    add_table_slide(
        prs,
        "Leave-One-Out Modality Ablation (v21 pipeline)",
        ["Configuration", "Channels Removed", "Token Acc.", "Seq. Acc.", "Features"],
        without_rows_detailed,
        highlight_rows=[0, 5],
        font_size=11,
        subtitle="Baseline: 80.69% (all modalities, v21). Removing environmental improves to 84.36% (+3.67%); removing color drops to 79.01% (-1.68%)."
    )

    # ========== SECTION: Model Performance Analysis ==========
    add_section_slide(prs, "Detailed Model Performance", "05")

    # ========== SLIDE: Command Confusion Matrix ==========
    cmd_cm = f"{base_path}/outputs/real_confusion_matrices/confusion_command.png"
    add_image_slide(
        prs,
        "G-code Command Classification",
        cmd_cm,
        "G0: 89%, G1: 80%, G53: 45% (rare command confusion).",
        image_width=7
    )

    # ========== SLIDE: Parameter Confusion Matrix ==========
    param_cm = f"{base_path}/outputs/real_confusion_matrices/confusion_parameter.png"
    add_image_slide(
        prs,
        "Parameter Type Classification",
        param_cm,
        "F: 100%, R: 100%, X: 94%, Z: 99%, Y: 71% (X/Y confusion in similar motions).",
        image_width=7
    )

    # ========== SLIDE: Digit Accuracy Heatmap ==========
    digit_heatmap = f"{base_path}/outputs/archive/2024-12_best_runs/sensor_multihead_v3/visualizations/paper/figures/digit_accuracy_heatmap.png"
    add_image_slide(
        prs,
        "Per-Digit Prediction Accuracy",
        digit_heatmap,
        "Accuracy degrades with digit precision: Sign (98-99%) -> D1/10s (94-96%) -> D6/.0001 (79-82%). Z-axis easiest.",
        image_width=7
    )

    # ========== SECTION: Architecture Ablations & Baselines ==========
    add_section_slide(prs, "Architecture Ablations and Baselines", "06")

    # ========== SLIDE: Architecture Ablations Table (FIX #4) ==========
    add_table_slide(
        prs,
        "Architecture and Hyperparameter Ablations (v21 pipeline)",
        ["Configuration", "Token Acc.", "Change vs Medium"],
        [
            ["Medium model (baseline)", "78.62%", "-"],
            ["XLarge model", "78.47%", "-0.15%"],
            ["Augmentation prob 0.5", "78.32%", "-0.30%"],
            ["Curriculum learning ON", "77.82%", "-0.80%"],
            ["No augmentation", "77.71%", "-0.91%"],
            ["No operation conditioning", "76.29%", "-2.33%"],
            ["Base model", "76.25%", "-2.37%"],
            ["No positional encoding", "76.18%", "-2.44%"],
            ["Tiny model", "72.97%", "-5.65%"],
            ["LSTM only (no transformer)", "70.82%", "-7.80%"],
        ],
        highlight_rows=[0, 9],
        font_size=12,
        subtitle="Run on v21 pipeline (baseline 78.62%). Model capacity (tiny to xlarge) changes accuracy by ~5%. Bottleneck is not model size."
    )

    # ========== SLIDE: Baselines ==========
    add_table_slide(
        prs,
        "Baseline Comparisons",
        ["Model", "Token Acc.", "Seq. Acc.", "Notes"],
        [
            ["Random guess", "0.5%", "0.0%", "1/2000 tokens"],
            ["Majority class", "23.74%", "0.0%", "Always predict most common"],
            ["LSTM only (no transformer)", "70.82%", "20.5%", "Encoder only (v21)"],
            ["Full model (all 16 sensors, v22)", "84.28%", "50.6%", "296 features"],
            ["Single best sensor (frame_r2)*", "93.50%", "77.1%", "41 features (v21)"],
        ],
        highlight_rows=[4],
        font_size=13,
        subtitle="*Single-sensor result from v21 pipeline. Clear progression from trivial baselines to full model."
    )

    # ========== SLIDE: Statistical Comparison ==========
    stat_comp = f"{base_path}/outputs/archive/2024-12_best_runs/sensor_multihead_v3/visualizations/paper/figures/statistical_comparison.png"
    add_image_slide(
        prs,
        "Statistical Comparison: Our Model vs Baselines",
        stat_comp,
        "Cross-validation score distributions (left) and Cohen's d effect sizes (right). All differences p < 0.001.",
        image_width=11
    )

    # ========== SECTION: Diagnosis & Conclusions ==========
    add_section_slide(prs, "Diagnosis and Conclusions", "07")

    # ========== SLIDE: Why Multi-Sensor Underperforms ==========
    add_content_slide(
        prs,
        "Why Does Multi-Sensor Fusion Underperform?",
        [
            "**1. Curse of Dimensionality + Non-Functional Sensor Noise**",
            "  - 41 features (single sensor): 3,160 training samples likely sufficient",
            "  - 296 features (all sensors): includes 68 features from 4 non-functional sensors (degraded/zero data)",
            "  - Sample density drops exponentially in high-dimensional space",
            "",
            "**2. Inter-Sensor Redundancy (Per-Modality Analysis)**",
            "  - Pressure, Temperature, Ax: Pearson r = 1.00 between frame sensor pairs",
            "  - RMS: r > 0.99 across multiple sensor locations",
            "  - Redundant features inflate dimensionality without adding discriminative information",
            "",
            "**3. Gradient Interference**",
            "  - Highly correlated input features produce conflicting gradient signals during backpropagation",
            "  - Self-attention weights are diluted across near-duplicate feature representations",
            "",
            "**4. Color Channel Paradox**",
            "  - R, G, B, A channels are 97%+ correlated with each other (intra-modality redundancy)",
            "  - Yet removing all color features drops token accuracy by 8.91%",
            "  - Interpretation: color provides cross-modal information absent from IMU/environmental sensors",
        ]
    )

    # ========== SLIDE: V22 Summary Table ==========
    add_table_slide(
        prs,
        "Summary: V21 vs V22 Results",
        ["Finding", "V21", "V22", "Change"],
        [
            ["Multi-sensor token acc", "80.69%", "84.28%", "+3.59%"],
            ["Multi-sensor seq acc", "35.4%", "50.6%", "+15.2%"],
            ["Single sensor token acc", "93.50%", "Pending re-run", "TBD"],
            ["Performance gap", "12.81%", "9.22%", "Narrowed"],
            ["Model variance (3 seeds)", "N/A", "+/-0.66%", "Stable"],
            ["Color contribution", "N/A", "+8.91%", "Essential"],
            ["Per-modality correlation", "Masked (avg)", "r=1.00", "Confirmed"],
        ],
        highlight_rows=[0, 1, 4, 5, 6],
        font_size=13
    )

    # ========== SLIDE: Key Conclusions ==========
    add_content_slide(
        prs,
        "Key Conclusions",
        [
            "**1. Single Sensor Achieves Best Performance (93.50% token accuracy)**",
            "  - Top three candidates: frame_r2 (93.50%), xa_motor (93.35%), frame_r1 (93.04%)",
            "  - Outperforms multi-sensor fusion (296 features) by 9.2 percentage points",
            "",
            "**2. V22 Normalization Fix Substantially Improved Multi-Sensor Results**",
            "  - Token accuracy: 80.69% -> 84.28% (+3.59 pp)",
            "  - Sequence accuracy: 35.4% -> 50.6% (+15.2 pp)",
            "",
            "**3. Color Sensors Contribute Essential Cross-Modal Information**",
            "  - Removing RGBA channels reduces token accuracy by 8.91 pp",
            "  - Despite 97%+ intra-channel correlation, color captures information absent from IMU sensors",
            "",
            "**4. Prof. Sodhi's Redundancy Hypothesis Confirmed**",
            "  - Frame sensors exhibit r = 1.00 for Pressure, Temperature, and Ax",
            "  - Original per-sensor averaging masked true per-modality redundancy",
            "",
            "**5. Architecture is Robust Across Random Seeds**",
            "  - Token accuracy variance: +/-0.66% across seeds {42, 123, 456}",
        ]
    )

    # ========== SLIDE: Recommended Path Forward (FIX #5) ==========
    add_two_column_slide(
        prs,
        "Recommended Path Forward",
        [
            "**Option A: Single Best Sensor**",
            "Current accuracy: 93.50% with 41 features",
            "Pros: Best performance, simplest, fastest",
            "Cons: May not generalize across machines",
            "",
            "**Option B: Top 2-3 Non-Redundant Sensors**",
            "Accuracy: TBD (experiments not yet run)",
            "Pros: More robust, complementary information",
            "Cons: Need to identify non-redundant pairs",
        ],
        [
            "**Option C: All Sensors + Feature Selection**",
            "Accuracy: TBD (PCA/selection not yet run)",
            "Pros: Uses all spatial positions",
            "Cons: Still more complex than needed",
            "",
            "**Recommended Next Steps**",
            "1. Test sensor combinations (e.g., frame_r2 + xa_motor)",
            "2. Remove redundant modalities (keep one of P/T/Ax)",
            "3. Try per-sensor encoders before fusion",
            "4. Collect more data if multi-sensor required",
        ]
    )

    # ========== SLIDE: Follow-Up Status (FIX #7, #8) ==========
    add_table_slide(
        prs,
        "Follow-Up Item Status",
        ["Item", "Status", "Result / Notes"],
        [
            ["Data leakage fix (train-only scaler)", "Complete", "+3.59% token accuracy"],
            ["Model stability (3 seeds)", "Complete", "+/-0.66% variance"],
            ["Color channel ablation", "Complete", "+8.91% contribution"],
            ["Per-modality correlation analysis", "Complete", "r=1.00 confirmed"],
            ["Single-sensor rankings (all 12)", "Complete", "frame_r2 best (93.50%)"],
            ["Architecture ablations (10 configs)", "Complete", "Capacity not the bottleneck"],
            ["Statistical comparison (Cohen's d)", "Complete", "All p < 0.001"],
            ["Top 2-3 sensor combinations", "Not started", "Planned"],
            ["PCA / feature selection (296 features)", "Not started", "Planned"],
            ["Embedding-level sensor fusion", "Not started", "Planned"],
            ["Per-sensor normalization comparison", "Not started", "Planned"],
            ["Top-k accuracy analysis", "Not started", "Planned"],
        ],
        highlight_rows=[7, 8, 9, 10, 11],
        font_size=10,
        subtitle="7 of 12 follow-up items completed. 5 remaining experiments planned."
    )

    # ========== SLIDE: Architecture Update Note ==========
    add_content_slide(
        prs,
        "Pending: Architecture Update",
        [
            "**Current Code**",
            "  - Unidirectional LSTM encoder",
            "  - Mean pooling over sequence",
            "",
            "**Proposed Update (per architecture diagrams)**",
            "  - Bidirectional LSTM encoder",
            "  - Attention pooling over sequence",
            "",
            "**Impact**",
            "  - If approved, all experiments must be re-run with updated architecture",
            "  - Sensor and modality ablations (currently v21) would also be updated to v22 pipeline",
            "  - Estimated scope: 71+ ablation runs + 3 ensemble seeds",
            "",
            "**Request**",
            "  - Group approval needed before committing to re-run",
        ]
    )

    # ========== SLIDE: Experiment Statistics (FIX #7) ==========
    add_two_column_slide(
        prs,
        "Experiment Statistics",
        [
            "**Total Experiments**",
            "290+ Weights & Biases runs",
            "71+ ablation studies",
            "~50 GPU hours training",
            "",
            "**Ablation Categories**",
            "12 single-sensor experiments",
            "12 leave-one-sensor-out",
            "6 single-modality experiments",
            "6 leave-one-modality-out",
            "35+ architecture ablations",
            "3 ensemble seed runs",
            "1 no-color ablation",
        ],
        [
            "**Pipeline Versions**",
            "v21: Original pipeline (sensor/modality/arch ablations)",
            "v22: Corrected normalization (multi-sensor, ensemble, color)",
            "Pending: Re-run all ablations on v22 pipeline",
            "",
            "**Reproducibility**",
            "Seeds: 42, 123, 456",
            "All runs logged to Weights & Biases",
            "Configuration files version-controlled",
        ]
    )

    # ========== SLIDE: Questions ==========
    add_title_slide(
        prs,
        "Questions?",
        "Stephen Eacuello  |  seacuello@uri.edu  |  January 2026"
    )

    # ================================================================
    # APPENDIX
    # ================================================================
    add_section_slide(prs, "Appendix", "A")

    # ========== APPENDIX: High Correlation Table ==========
    corr_rows = [[mod, pair, r] for mod, pair, r in HIGH_CORRELATIONS]
    add_table_slide(
        prs,
        "Appendix: High Per-Modality Correlations (|r| > 0.97)",
        ["Modality", "Sensor Pair", "Correlation"],
        corr_rows,
        highlight_rows=[0, 1, 2, 3],
        font_size=12,
        subtitle="Environmental and RMS modalities show near-perfect correlation across frame sensors."
    )

    # ========== APPENDIX: Sensor Location Importance ==========
    loc_bars = f"{base_path}/outputs/archive/2024-12_best_runs/sensor_multihead_v3/visualizations/paper/figures/location_importance_bars.png"
    add_image_slide(
        prs,
        "Appendix: Sensor Location Importance (Leave-One-Out)",
        loc_bars,
        "Accuracy drop when each sensor location is removed. y_bed_3 and frame_r1 are most important for multi-sensor model.",
        image_width=8
    )

    # ========== APPENDIX: Modality x Location Heatmap ==========
    loc_heatmap = f"{base_path}/outputs/archive/2024-12_best_runs/sensor_multihead_v3/visualizations/paper/figures/sensor_location_heatmap.png"
    add_image_slide(
        prs,
        "Appendix: Modality x Location Interaction Matrix",
        loc_heatmap,
        "Most critical: Proximity at frame_r1 (-2.77% accuracy drop). Color and Motor Current contribute least.",
        image_width=10
    )

    # ========== APPENDIX: Modality Correlation Heatmap ==========
    mod_corr = f"{base_path}/outputs/analysis/modality_correlation_heatmap.png"
    add_image_slide(
        prs,
        "Appendix: Modality Correlation Matrix",
        mod_corr,
        "Color channels R, G, B, A are 97%+ correlated. Ax-Gz, Gz-Pressure also highly correlated.",
        image_width=9
    )

    # Save presentation
    output_path = f"{base_path}/outputs/presentations/gcode_fingerprinting_v22_results_v4.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    return output_path


if __name__ == "__main__":
    create_presentation()
