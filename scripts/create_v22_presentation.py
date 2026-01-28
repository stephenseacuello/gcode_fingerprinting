#!/usr/bin/env python3
"""
Create G-code Fingerprinting v22 Results PowerPoint Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os

# Constants
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 112, 192)
GREEN = RGBColor(0, 128, 0)
RED = RGBColor(192, 0, 0)
ORANGE = RGBColor(255, 153, 0)
GRAY = RGBColor(89, 89, 89)


def set_cell_text(cell, text, font_size=14, bold=False, color=None):
    """Set text in a table cell with formatting."""
    cell.text = str(text)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        if color:
            paragraph.font.color.rgb = color


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.33), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullet_points=None, has_table=False):
    """Add a content slide with title and optional bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title bar background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Bullet points
    if bullet_points:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.33), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Check if it's a header (starts with **)
            if point.startswith("**") and point.endswith("**"):
                p.text = point.strip("*")
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = DARK_BLUE
                p.space_before = Pt(18)
            elif point.startswith("  - "):
                p.text = point.strip()
                p.font.size = Pt(18)
                p.level = 1
            else:
                p.text = f"• {point}"
                p.font.size = Pt(20)
                p.space_before = Pt(6)

    return slide


def add_table_slide(prs, title, headers, rows, highlight_rows=None):
    """Add a slide with a table."""
    slide = add_content_slide(prs, title)

    # Create table
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.33)
    height = Inches(0.4 * num_rows)

    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    # Set column widths
    col_width = width / num_cols
    for i in range(num_cols):
        table.columns[i].width = int(col_width)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        set_cell_text(cell, header, font_size=16, bold=True)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # Data rows
    highlight_rows = highlight_rows or []
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            is_highlight = row_idx in highlight_rows
            set_cell_text(cell, cell_data, font_size=14, bold=is_highlight)
            if is_highlight:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 242, 204)

    return slide


def add_two_column_slide(prs, title, left_content, right_content):
    """Add a slide with two columns."""
    slide = add_content_slide(prs, title)

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.9), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if point.startswith("**"):
            p.text = point.strip("*")
            p.font.bold = True
            p.font.size = Pt(20)
            p.font.color.rgb = DARK_BLUE
        else:
            p.text = f"• {point}"
            p.font.size = Pt(16)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.4), Inches(5.9), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if point.startswith("**"):
            p.text = point.strip("*")
            p.font.bold = True
            p.font.size = Pt(20)
            p.font.color.rgb = DARK_BLUE
        else:
            p.text = f"• {point}"
            p.font.size = Pt(16)

    return slide


def add_image_slide(prs, title, image_path, caption=""):
    """Add a slide with an image."""
    slide = add_content_slide(prs, title)

    if os.path.exists(image_path):
        # Add image centered
        img_left = Inches(1.5)
        img_top = Inches(1.5)
        img_width = Inches(10)
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)

    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    return slide


def create_presentation():
    """Create the full v22 presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ========== SLIDE 1: Title ==========
    add_title_slide(
        prs,
        "G-code Fingerprinting:\n2K Vocab Experiment Results",
        "Stephen Eacuello | January 2026 | Version 22"
    )

    # ========== SLIDE 2: Executive Summary ==========
    slide = add_table_slide(
        prs,
        "Executive Summary: Key Results (v22)",
        ["Configuration", "Token Accuracy", "Sequence Accuracy", "Features"],
        [
            ["Best Single Sensor (frame_r2)", "93.50%", "77.1%", "41"],
            ["All 12 Sensors (v22 corrected)", "84.28%", "50.6%", "296"],
            ["All 12 Sensors (without color)", "75.37%", "28.0%", "208"],
            ["Performance Gap", "-9.22%", "-26.5%", "+255"],
        ],
        highlight_rows=[0, 3]
    )

    # Add key findings text below table
    findings_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.33), Inches(2.5))
    tf = findings_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key v22 Changes:"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_BLUE

    findings = [
        "Fixed data leakage (scaler fitted on training only): +3.59% token accuracy",
        "Color sensors essential: +8.91% contribution",
        "Model stable across seeds: 83.67% ± 0.66%"
    ]
    for finding in findings:
        p = tf.add_paragraph()
        p.text = f"• {finding}"
        p.font.size = Pt(18)

    # ========== SLIDE 3: Dataset Overview ==========
    add_two_column_slide(
        prs,
        "Dataset Overview",
        [
            "**Data Composition**",
            "145 aligned CSV files from real CNC operations",
            "3,160+ training sequences (64 timesteps each)",
            "2,000+ token vocabulary (4-digit precision)",
            "",
            "**Operations**",
            "Face operations: 40 files (27.6%)",
            "Pocket operations: 40 files (27.6%)",
            "Adaptive operations: 40 files (27.6%)",
            "Damaged operations: 15 files (10.3%)",
        ],
        [
            "**Sensors (12 total)**",
            "Frame: frame_r1, frame_r2, frame_b2, frame_l2, frame_l3",
            "Spindle: spindle1, spindle2",
            "Y-bed: y_bed__1, y_bed__2, y_bed__3, y_bed__4",
            "Motor: xa_motor",
            "",
            "**Modalities per sensor (17)**",
            "Motion: Ax, Ay, Az, Gx, Gy, Gz, Mx, My, Mz",
            "Environmental: Pressure, Temperature, Proximity",
            "Color: ColorR, ColorG, ColorB, ColorA",
            "Computed: RMS",
        ]
    )

    # ========== SLIDE 4: Two-Stage Architecture ==========
    add_content_slide(
        prs,
        "Two-Stage Architecture",
        [
            "**Stage 1: Encoder (MM-DTAE-LSTM) - Frozen after pre-training**",
            "  - Input: Sensor data [B, 64, 296]",
            "  - Output: Memory embeddings [B, 64, 128]",
            "  - Achieves 100% operation classification",
            "",
            "**Stage 2: Decoder (Transformer Multi-Head)**",
            "  - 4 layers, 8 attention heads, d_model=192",
            "  - Multi-head output: Type, Command, Param Type, Digit Values",
            "",
            "**Training Configuration**",
            "  - Batch size: 32, Learning rate: 0.0002 (AdamW)",
            "  - Max epochs: 150 (patience: 30)",
            "  - Seeds tested: 42, 123, 456",
        ]
    )

    # ========== SLIDE 5: Data Leakage Fix (NEW) ==========
    slide = add_table_slide(
        prs,
        "v22 Improvement: Data Leakage Fix",
        ["Metric", "v21 (With Leakage)", "v22 (Corrected)", "Change"],
        [
            ["Multi-sensor token acc", "80.69%", "84.28%", "+3.59%"],
            ["Multi-sensor seq acc", "35.4%", "50.6%", "+15.2%"],
            ["Single sensor token acc", "93.50%", "93.50%", "unchanged"],
            ["Single sensor seq acc", "77.1%", "77.1%", "unchanged"],
        ],
        highlight_rows=[0, 1]
    )

    # Add explanation
    exp_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.33), Inches(2.5))
    tf = exp_box.text_frame
    tf.word_wrap = True
    explanations = [
        ("Issue:", "StandardScaler was fitted on entire dataset before train/val/test split"),
        ("Fix:", "StandardScaler now fitted ONLY on training data"),
        ("Impact:", "Multi-sensor model benefits more from proper methodology"),
    ]
    for i, (label, text) in enumerate(explanations):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{label} {text}"
        p.font.size = Pt(16)

    # ========== SLIDE 6: Model Stability (NEW) ==========
    slide = add_table_slide(
        prs,
        "v22 Improvement: Model Stability Across Seeds",
        ["Seed", "Token Acc.", "Seq. Acc.", "Val. Token"],
        [
            ["42", "84.28%", "50.57%", "85.83%"],
            ["123", "82.98%", "50.08%", "84.76%"],
            ["456", "83.75%", "47.80%", "85.37%"],
            ["Mean", "83.67%", "49.48%", "85.32%"],
            ["Std", "±0.66%", "±1.52%", "±0.54%"],
        ],
        highlight_rows=[3, 4]
    )

    # Add takeaways
    take_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.33), Inches(2))
    tf = take_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Takeaways:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_BLUE

    takeaways = [
        "Very low variance (±0.66%) demonstrates robust training",
        "Results are independent of random initialization",
        "Model architecture is well-tuned"
    ]
    for t in takeaways:
        p = tf.add_paragraph()
        p.text = f"• {t}"
        p.font.size = Pt(16)

    # ========== SLIDE 7: Color Channel Ablation (NEW) ==========
    slide = add_table_slide(
        prs,
        "v22 Improvement: Color Sensors Are Essential",
        ["Configuration", "Features", "Token Acc.", "Seq. Acc."],
        [
            ["With Color (RGBA)", "296", "84.28%", "50.57%"],
            ["Without Color", "208", "75.37%", "27.97%"],
            ["Difference", "-88", "-8.91%", "-22.60%"],
        ],
        highlight_rows=[2]
    )

    # Add explanation
    exp_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.33), Inches(3))
    tf = exp_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Why Color Matters - captures ambient light changes from:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_BLUE

    reasons = [
        "Chip formation dynamics - cutting creates debris patterns",
        "Coolant behavior - liquid flow changes light reflection",
        "Surface reflectance - material removal changes surface",
        "Tool engagement - tool-workpiece interaction affects lighting",
    ]
    for r in reasons:
        p = tf.add_paragraph()
        p.text = f"• {r}"
        p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Surprising: Despite R,G,B,A being 97%+ correlated internally, removing them drops accuracy by 8.9%!"
    p.font.italic = True
    p.font.size = Pt(16)
    p.font.color.rgb = RED

    # ========== SLIDE 8: Single Sensor Performance ==========
    add_table_slide(
        prs,
        "Single Sensor Performance Ranking",
        ["Rank", "Sensor", "Token Acc.", "Seq Acc.", "Features"],
        [
            ["1", "frame_r2", "93.50%", "77.1%", "41"],
            ["2", "xa_motor", "93.35%", "76.6%", "41"],
            ["3", "frame_r1", "93.04%", "74.6%", "41"],
            ["4", "y_bed__1", "92.93%", "74.8%", "41"],
            ["5", "spindle2", "92.54%", "73.0%", "41"],
            ["...", "...", "...", "...", "..."],
            ["12", "spindle1", "88.07%", "58.5%", "41"],
            ["-", "All 12 Sensors (v22)", "84.28%", "50.6%", "296"],
        ],
        highlight_rows=[0, 7]
    )

    # ========== SLIDE 9: Per-Modality Correlation (NEW) ==========
    add_table_slide(
        prs,
        "Per-Modality Correlation: Frame Sensors ARE Correlated!",
        ["Modality", "Sensor Pair", "Correlation"],
        [
            ["Pressure", "frame_r1 ↔ frame_b2", "1.000"],
            ["Temperature", "frame_r1 ↔ frame_b2", "1.000"],
            ["Ax", "frame_r1 ↔ frame_b2", "1.000"],
            ["RMS", "frame_r2 ↔ y_bed__3", "0.997"],
            ["My", "frame_r1 ↔ frame_b2", "-0.990"],
            ["Mz", "frame_r1 ↔ frame_b2", "0.988"],
            ["Gz", "frame_r1 ↔ frame_b2", "0.973"],
        ],
        highlight_rows=[0, 1, 2]
    )

    # ========== SLIDE 10: Per-Modality Correlation Figure ==========
    image_path = "outputs/jan23_followup/correlations/per_modality_correlation_summary.png"
    add_image_slide(
        prs,
        "Per-Modality Sensor Correlations (Visual)",
        image_path,
        "Professor Sodhi's intuition confirmed: Adjacent frame sensors show near-perfect correlation when viewed per-modality"
    )

    # ========== SLIDE 11: Single Modality Performance ==========
    add_table_slide(
        prs,
        "Single Modality Performance",
        ["Modality", "Token Acc.", "Seq Acc.", "Features"],
        [
            ["RMS only", "90.44%", "67.5%", "40"],
            ["Gyroscope only", "89.06%", "61.3%", "72"],
            ["Magnetometer only", "87.92%", "58.2%", "72"],
            ["Color only", "86.69%", "55.6%", "88"],
            ["Accelerometer only", "86.50%", "51.7%", "72"],
            ["Environmental only", "82.72%", "39.7%", "72"],
            ["All Modalities (v22)", "84.28%", "50.6%", "296"],
        ],
        highlight_rows=[0, 6]
    )

    # ========== SLIDE 12: Modality Correlation ==========
    add_table_slide(
        prs,
        "Modality Correlation Analysis",
        ["Pair", "Correlation", "Implication"],
        [
            ["ColorG ↔ ColorB", "0.998", "Color channels nearly identical"],
            ["ColorR ↔ ColorB", "0.994", "Color channels nearly identical"],
            ["ColorR ↔ ColorG", "0.992", "Color channels nearly identical"],
            ["Ax ↔ Gz", "-0.877", "Motion coupling"],
            ["Ax ↔ Pressure", "-0.861", "Physics relationship"],
            ["Gz ↔ Pressure", "0.843", "Physics relationship"],
        ],
        highlight_rows=[0, 1, 2]
    )

    # ========== SLIDE 13: Architecture Ablations ==========
    add_table_slide(
        prs,
        "Model Size Has Minimal Impact",
        ["Configuration", "Token Accuracy", "Change"],
        [
            ["Medium model", "78.62%", "baseline"],
            ["XLarge model", "78.47%", "-0.15%"],
            ["Aug prob 0.5", "78.32%", "-0.30%"],
            ["Curriculum ON", "77.82%", "-0.80%"],
            ["Base model", "76.25%", "-2.37%"],
            ["Tiny model", "72.97%", "-5.65%"],
        ],
        highlight_rows=[0]
    )

    # ========== SLIDE 14: Baseline Comparisons ==========
    add_table_slide(
        prs,
        "Baseline Comparisons",
        ["Model", "Token Acc.", "Seq Acc."],
        [
            ["Random guess", "0.5%", "0.0%"],
            ["Majority class", "23.74%", "0.0%"],
            ["LSTM only", "70.82%", "20.5%"],
            ["Full model (all sensors, v22)", "84.28%", "50.6%"],
            ["Single best sensor (frame_r2)", "93.50%", "77.1%"],
        ],
        highlight_rows=[4]
    )

    # ========== SLIDE 15: Why Does This Happen? ==========
    add_content_slide(
        prs,
        "Why Does Multi-Sensor Fusion Underperform?",
        [
            "**1. Sample Complexity**",
            "  - With 41 features: 3,160 samples may be sufficient",
            "  - With 296 features: 3,160 samples is sparse in feature space",
            "",
            "**2. Redundant Information**",
            "  - Per-modality analysis shows sensors ARE highly correlated",
            "  - Environmental/RMS modalities nearly identical across frame sensors",
            "",
            "**3. Feature Interference**",
            "  - Correlated features provide conflicting gradients",
            "  - Attention mechanism diluted across many similar inputs",
            "",
            "**4. The Color Paradox**",
            "  - Color R,G,B,A are 97%+ correlated (redundant within modality)",
            "  - Yet removing color drops accuracy by 8.9%",
            "  - Color provides CROSS-MODAL information not in other sensors",
        ]
    )

    # ========== SLIDE 16: v22 Summary ==========
    slide = add_table_slide(
        prs,
        "Summary: v21 vs v22 Results",
        ["Finding", "v21", "v22", "Change"],
        [
            ["Multi-sensor token acc", "80.69%", "84.28%", "+3.59%"],
            ["Multi-sensor seq acc", "35.4%", "50.6%", "+15.2%"],
            ["Single sensor token acc", "93.50%", "93.50%", "unchanged"],
            ["Performance gap", "12.81%", "9.22%", "narrowed"],
            ["Model variance", "N/A", "±0.66%", "stable"],
            ["Color contribution", "N/A", "+8.91%", "essential"],
        ],
        highlight_rows=[0, 1, 4, 5]
    )

    # ========== SLIDE 17: Conclusions ==========
    add_content_slide(
        prs,
        "Updated Conclusions (v22)",
        [
            "**1. Single sensor still achieves best performance (93.50%)**",
            "  - frame_r2, xa_motor, and frame_r1 all viable choices",
            "",
            "**2. Multi-sensor gap narrowed after fixing methodology**",
            "  - From ~13% to ~9% gap after data leakage fix",
            "  - Sequence accuracy improved dramatically: 35.4% → 50.6%",
            "",
            "**3. Color sensors are essential despite high internal correlation**",
            "  - Removing color drops accuracy by 8.91%",
            "  - Provides unique cross-modal information",
            "",
            "**4. Professor Sodhi's intuition confirmed**",
            "  - Frame sensors show r=1.00 correlation for Pressure, Temperature, Ax",
            "  - Original analysis masked this by averaging modalities",
            "",
            "**5. Model is stable across random seeds (±0.66%)**",
        ]
    )

    # ========== SLIDE 18: Next Steps ==========
    add_two_column_slide(
        prs,
        "Proposed Next Steps",
        [
            "**Completed ✓**",
            "Fixed data leakage (train-only scaler)",
            "Validated model stability (3 seeds)",
            "Color channel ablation",
            "Per-modality correlation analysis",
            "",
            "**Priority 1: Further Validation**",
            "Run single sensor with corrected methodology",
            "Cross-validate color findings",
            "Generate publication-quality figures",
        ],
        [
            "**Priority 2: Sensor Optimization**",
            "Test 2-3 sensor combinations",
            "Test removing redundant sensors",
            "PCA on redundant modalities",
            "",
            "**Priority 3: Architecture**",
            "Sensor-specific attention weights",
            "Hierarchical fusion",
            "Modality-wise dropout",
        ]
    )

    # ========== SLIDE 19: Recommended Path Forward ==========
    add_content_slide(
        prs,
        "Recommended Path Forward",
        [
            "**Option A: Use Single Best Sensor (Simplest)**",
            "  - Accuracy: 93.50% | Features: 41",
            "  - Pros: Best performance, simplest model, fastest inference",
            "  - Cons: May not generalize to other machines",
            "",
            "**Option B: Top 2-3 Non-Redundant Sensors**",
            "  - Expected: 90-92% | Features: 82-123",
            "  - Pros: More robust, uses complementary information",
            "  - Cons: Need to identify truly non-redundant sensors",
            "",
            "**Option C: All Sensors with Modality Selection**",
            "  - Expected: 86-88% | Features: ~150",
            "  - Pros: Uses all spatial positions",
            "  - Cons: Still more complex than needed",
        ]
    )

    # ========== SLIDE 20: Experiment Statistics ==========
    add_two_column_slide(
        prs,
        "Experiment Statistics",
        [
            "**Compute Resources**",
            "Total W&B runs: 290+",
            "Total training time: ~50+ GPU hours",
            "Storage: ~1.2GB checkpoints",
            "",
            "**Reproducibility**",
            "All experiments logged to W&B",
            "Seeds: 42, 123, 456",
            "Config: configs/best_lambda_2k_vocab.json",
        ],
        [
            "**Key v22 Checkpoints**",
            "Best multi-sensor: jan23_followup/no_leakage/training_tuned/",
            "Ensemble: jan23_followup/ensemble/seed_{42,123,456}/",
            "No-color ablation: jan23_followup/no_color/training_tuned/",
            "Best single sensor: full_pipeline_2k_vocab/sensor_ablations/only_frame_r2/",
        ]
    )

    # ========== SLIDE 21: Thank You ==========
    add_title_slide(
        prs,
        "Questions?",
        "Stephen Eacuello | seacuello@uri.edu | January 2026"
    )

    # Save presentation
    output_path = "outputs/presentations/gcode_fingerprinting_v22_results.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
