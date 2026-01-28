#!/usr/bin/env python3
"""
Create PowerPoint presentation for G-code Fingerprinting 2K Vocab Results.

Usage:
    python scripts/create_presentation.py

Author: Claude Code
Date: January 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 112, 192)
GREEN = RGBColor(0, 128, 0)
RED = RGBColor(192, 0, 0)
ORANGE = RGBColor(255, 153, 0)
GRAY = RGBColor(89, 89, 89)


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, title):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(10), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, sub_bullets=None):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)

        # Add sub-bullets if provided
        if sub_bullets and i in sub_bullets:
            for sub in sub_bullets[i]:
                sp = tf.add_paragraph()
                sp.text = sub
                sp.font.size = Pt(16)
                sp.font.color.rgb = GRAY
                sp.level = 1
                sp.space_before = Pt(3)

    return slide


def add_table_slide(prs, title, headers, rows, highlight_rows=None):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Table
    n_rows = len(rows) + 1
    n_cols = len(headers)

    table_width = Inches(9)
    table_height = Inches(0.4 * n_rows)
    left = Inches(0.5)
    top = Inches(1.3)

    table = slide.shapes.add_table(n_rows, n_cols, left, top, table_width, table_height).table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(14)
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.alignment = PP_ALIGN.CENTER

            # Highlight specific rows
            if highlight_rows and row_idx in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(198, 224, 180)
                para.font.bold = True

    return slide


def add_key_finding_slide(prs, title, finding, details):
    """Add a slide highlighting a key finding."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Key finding box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 112, 192)

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = finding
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Details
    detail_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(3.5))
    tf = detail_box.text_frame
    tf.word_wrap = True

    for i, detail in enumerate(details):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = detail
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        p.space_before = Pt(8)

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "G-code Fingerprinting",
        "2K Vocabulary Experiment Results\n\nStephen Eacuello | January 2026\nELE 588 Applied Machine Learning"
    )

    # Slide 2: Executive Summary
    add_key_finding_slide(
        prs,
        "Executive Summary: Curse of Dimensionality Confirmed",
        "Single Sensor: 93.5%  vs  All 12 Sensors: 80.7%  =  -12.8% Gap",
        [
            "Best single sensor (frame_r2): 93.50% accuracy with 41 features",
            "All 12 sensors combined: 80.69% accuracy with 296 features",
            "Every single sensor individually outperforms the combined model",
            "This is NOT a redundancy issue - sensors have low correlation (max |r| = 0.56)",
            "Bottom line: More features = worse performance due to dimensionality"
        ]
    )

    # Slide 3: Dataset Overview
    add_content_slide(
        prs,
        "Dataset Overview",
        [
            "Data: 145 aligned CSV files from real CNC operations",
            "Sequences: 3,160+ training sequences (64 timesteps each)",
            "Vocabulary: 2,000+ tokens (4-digit precision)",
            "Operations: Face (40), Pocket (40), Adaptive (40), Damaged (15)",
            "Sensors: 12 total - Frame (5), Spindle (2), Y-bed (4), Motor (1)",
            "Modalities per sensor: 17 - Motion (9), Environmental (3), Color (4), RMS (1)",
            "Total features: 12 sensors × 17 modalities ≈ 296 features"
        ]
    )

    # Slide 4: Experimental Setup
    add_content_slide(
        prs,
        "Experimental Setup",
        [
            "Two-Stage Architecture:",
            "Training Configuration:",
            "Experiments Run: 71+ ablation studies",
            "W&B Tracking: 290 runs logged"
        ],
        {
            0: ["Stage 1: MM-DTAE-LSTM Encoder (frozen, 100% op classification)",
                "Stage 2: Transformer Multi-Head Decoder (4 layers, 8 heads)"],
            1: ["Batch size: 32, Learning rate: 0.0002 (AdamW)",
                "Max epochs: 100 (patience: 15), Seeds: 42, 123, 456"],
            2: ["12 single-sensor experiments",
                "12 leave-one-sensor-out experiments",
                "6 single-modality + 6 leave-one-modality experiments",
                "35+ architecture/hyperparameter ablations"]
        }
    )

    # Slide 5: Section - Sensor Analysis
    add_section_slide(prs, "Part 1: Single Sensor Analysis")

    # Slide 6: Sensor Ranking Table
    add_table_slide(
        prs,
        "Single Sensor Performance Ranking",
        ["Rank", "Sensor", "Test Accuracy", "Seq Accuracy", "Features"],
        [
            ["1", "frame_r2", "93.50%", "77.1%", "41"],
            ["2", "xa_motor", "93.35%", "76.6%", "41"],
            ["3", "frame_r1", "93.04%", "74.6%", "41"],
            ["4", "y_bed__1", "92.93%", "74.8%", "41"],
            ["5", "spindle2", "92.54%", "73.0%", "41"],
            ["6", "frame_l2", "92.35%", "72.0%", "41"],
            ["...", "...", "...", "...", "..."],
            ["12", "spindle1", "88.07%", "58.5%", "41"],
            ["-", "All 12 Sensors", "80.69%", "35.4%", "296"],
        ],
        highlight_rows=[0, 8]
    )

    # Slide 7: Sensor Sensitivity Analysis
    add_key_finding_slide(
        prs,
        "Sensor Sensitivity Analysis",
        "Top 3 Sensors Are Very Close in Performance",
        [
            "frame_r2: 93.50% - Best overall",
            "xa_motor: 93.35% - Only 0.15% behind",
            "frame_r1: 93.04% - Strong third option",
            "",
            "Pattern by sensor type:",
            "• Frame sensors (r1, r2, l2, l3, b2): 92.2% - 93.5% (all strong)",
            "• Y-bed sensors: 91.1% - 92.9% (consistent mid-tier)",
            "• Spindle: spindle2 (92.5%) >> spindle1 (88.1%)",
            "• Motor: 93.35% (excellent)"
        ]
    )

    # Slide 8: Sensor Correlation
    add_key_finding_slide(
        prs,
        "Sensor Correlation Analysis",
        "Sensors Are NOT Redundant (max |r| = 0.56)",
        [
            "Maximum correlation between any sensor pair: 0.56",
            "Mean absolute correlation: 0.17",
            "Pairs with |r| > 0.5: Only 1 (frame_r1 ↔ frame_b2: -0.56)",
            "Pairs with |r| > 0.8: None",
            "",
            "Interpretation:",
            "• Each sensor provides unique, non-redundant information",
            "• The performance gap is NOT from redundant sensors",
            "• The curse of dimensionality is from too many features to learn"
        ]
    )

    # Slide 9: Section - Modality Analysis
    add_section_slide(prs, "Part 2: Modality Analysis")

    # Slide 10: Modality Performance Table
    add_table_slide(
        prs,
        "Single Modality Performance",
        ["Modality", "Test Accuracy", "Seq Accuracy", "Features"],
        [
            ["RMS only", "90.44%", "67.5%", "40"],
            ["Gyroscope only", "89.06%", "61.3%", "72"],
            ["Magnetometer only", "87.92%", "58.2%", "72"],
            ["Color only", "86.69%", "55.6%", "88"],
            ["Accelerometer only", "86.50%", "51.7%", "72"],
            ["Environmental only", "82.72%", "39.7%", "72"],
            ["All Modalities", "80.69%", "35.4%", "296"],
        ],
        highlight_rows=[0, 6]
    )

    # Slide 11: Modality Correlation
    add_key_finding_slide(
        prs,
        "Modality Correlation Analysis",
        "Color Channels Are Highly Redundant (r > 0.97)",
        [
            "ColorG ↔ ColorB: r = 0.998",
            "ColorR ↔ ColorB: r = 0.994",
            "ColorR ↔ ColorG: r = 0.992",
            "ColorR ↔ ColorA: r = 0.987",
            "",
            "Other notable correlations:",
            "• Ax ↔ Gz: r = -0.877",
            "• Ax ↔ Pressure: r = -0.861",
            "",
            "Recommendation: Could reduce 4 color features to 1"
        ]
    )

    # Slide 12: Section - Architecture Analysis
    add_section_slide(prs, "Part 3: Architecture & Ablations")

    # Slide 13: Model Size Impact
    add_table_slide(
        prs,
        "Model Size Has Minimal Impact",
        ["Configuration", "Test Accuracy", "Change"],
        [
            ["Medium model", "78.62%", "baseline"],
            ["XLarge model", "78.47%", "-0.15%"],
            ["Aug prob 0.5", "78.32%", "-0.30%"],
            ["Curriculum ON", "77.82%", "-0.80%"],
            ["Base model", "76.25%", "-2.37%"],
            ["Tiny model", "72.97%", "-5.65%"],
        ],
        highlight_rows=[]
    )

    # Slide 14: Component Ablations
    add_table_slide(
        prs,
        "Component Ablations & Baselines",
        ["Configuration", "Token Acc", "Seq Acc"],
        [
            ["Random guess", "0.5%", "0.0%"],
            ["Majority class", "23.74%", "0.0%"],
            ["LSTM only (no transformer)", "70.82%", "20.5%"],
            ["No augmentation", "77.71%", "-"],
            ["No operation conditioning", "76.29%", "-"],
            ["No positional encoding", "76.18%", "-"],
            ["Full model (all sensors)", "80.69%", "35.4%"],
            ["Single best sensor", "93.50%", "77.1%"],
        ],
        highlight_rows=[7]
    )

    # Slide 15: Section - Diagnosis
    add_section_slide(prs, "Part 4: Diagnosing the Problem")

    # Slide 16: Training vs Validation
    add_key_finding_slide(
        prs,
        "Training vs Validation Analysis",
        "Combined Model Shows Mild Overfitting",
        [
            "Single sensor (frame_r2):",
            "• Best validation accuracy: 92.98%",
            "• Test accuracy: 93.50%",
            "• No overfitting (test > val)",
            "",
            "All 12 sensors:",
            "• Best validation accuracy: 82.84%",
            "• Test accuracy: 80.69%",
            "• Possible slight overfitting (val > test by ~2%)",
            "",
            "Key insight: The gap exists even in validation (82.84% vs 92.98%)"
        ]
    )

    # Slide 17: Why This Happens
    add_content_slide(
        prs,
        "Why Does This Happen?",
        [
            "Sample Complexity:",
            "Optimization Difficulty:",
            "Feature Interference:",
            "NOT Redundancy:"
        ],
        {
            0: ["With 41 features: 3,160 samples may be sufficient",
                "With 296 features: 3,160 samples is sparse in feature space",
                "Need exponentially more data as dimensions increase"],
            1: ["296 dimensions = harder loss landscape",
                "More local minima to get stuck in",
                "Gradient signal diluted across many features"],
            2: ["Irrelevant features add noise",
                "Model struggles to identify signal",
                "Regularization can't fully compensate"],
            3: ["Correlation analysis shows sensors are unique",
                "Problem is quantity, not quality of features"]
        }
    )

    # Slide 18: Section - Solutions
    add_section_slide(prs, "Part 5: Proposed Solutions")

    # Slide 19: Solution 1 - PCA
    add_content_slide(
        prs,
        "Solution 1: PCA / Feature Selection",
        [
            "Approach: Apply PCA to reduce 296 → 50 principal components",
            "Test variance thresholds: 90%, 95%, 99%",
            "Compare to best single sensor (41 features)",
            "",
            "Expected outcomes:",
            "• If PCA matches single sensor: Confirms we just need fewer dimensions",
            "• If PCA underperforms: Suggests sensor-specific information matters",
            "",
            "Status: To be implemented"
        ]
    )

    # Slide 20: Solution 2 & 3
    add_content_slide(
        prs,
        "Solutions 2 & 3: Regularization & Sensor Fusion",
        [
            "Regularization (already tested):",
            "",
            "Sensor Fusion Alternatives:",
            "",
            "Sensor Combination Testing:"
        ],
        {
            0: ["Dropout tested (0.1-0.5): minimal impact (~2% range)",
                "L1/L2 regularization: not yet tested"],
            2: ["Option A: Per-sensor encoders (12 branches, fuse at embedding)",
                "Option B: Attention-based fusion (learn sensor importance)",
                "Option C: Raw concatenation (current approach)"],
            4: ["Test top 2-sensor combinations",
                "Test top 3-sensor combinations",
                "Find diminishing returns threshold"]
        }
    )

    # Slide 21: Section - Next Steps
    add_section_slide(prs, "Part 6: Recommendations")

    # Slide 22: Recommended Path
    add_content_slide(
        prs,
        "Recommended Path Forward",
        [
            "Option A: Use Single Best Sensor (Simplest)",
            "",
            "Option B: Top 2-3 Sensor Ensemble",
            "",
            "Option C: PCA-Reduced Full Data"
        ],
        {
            0: ["Accuracy: 93.50% | Features: 41",
                "Pros: Best performance, simplest model",
                "Cons: May not generalize to other machines"],
            2: ["Expected: 90-92% | Features: 82-123",
                "Pros: More robust, still good performance",
                "Cons: Need to validate combinations don't hurt"],
            4: ["Expected: 85-88% | Features: 50 (from 296)",
                "Pros: Uses all sensor information",
                "Cons: May lose sensor-specific patterns"]
        }
    )

    # Slide 23: Immediate Actions
    add_content_slide(
        prs,
        "Immediate Next Steps",
        [
            "Priority 1: Validate Findings",
            "",
            "Priority 2: Test Professor's Solutions",
            "",
            "Priority 3: Sensor Combinations"
        ],
        {
            0: ["Run top 3 sensors with different seeds (42, 123, 456)",
                "Compute training accuracy curves for single vs combined",
                "Generate top-5 and top-10 accuracy metrics"],
            2: ["Implement PCA reduction (296 → 50)",
                "Sweep L1/L2 regularization",
                "Test per-sensor normalization"],
            4: ["Test 2-sensor combinations (top 5 pairs)",
                "Test 3-sensor combinations (top 5 triples)",
                "Find diminishing returns threshold"]
        }
    )

    # Slide 24: Summary Statistics
    add_table_slide(
        prs,
        "Experiment Summary Statistics",
        ["Category", "# Runs", "Key Finding"],
        [
            ["Single Sensor", "12", "All beat combined (88-93%)"],
            ["Single Modality", "6", "RMS alone = 90.4%"],
            ["Leave-One-Out Sensor", "12", "Removing any sensor helps slightly"],
            ["Leave-One-Out Modality", "6", "Removing environmental helps most"],
            ["Architecture Ablations", "35+", "Model capacity not the bottleneck"],
            ["Total", "71+", "Curse of dimensionality confirmed"],
        ],
        highlight_rows=[5]
    )

    # Slide 25: Questions
    add_title_slide(
        prs,
        "Questions?",
        "Stephen Eacuello | seacuello@uri.edu\n\nAll experiments logged to W&B (290 runs)\nResults: outputs/full_pipeline_2k_vocab/"
    )

    # Save presentation
    output_path = "outputs/presentations/gcode_fingerprinting_2k_vocab_results.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    main()
