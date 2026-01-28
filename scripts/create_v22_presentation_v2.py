#!/usr/bin/env python3
"""
Create G-code Fingerprinting v22 Results PowerPoint Presentation - V2 (Polished)
Includes complete tables, confusion matrices, and all figures.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Constants
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 112, 192)
GREEN = RGBColor(0, 128, 0)
RED = RGBColor(192, 0, 0)
ORANGE = RGBColor(255, 153, 0)
GRAY = RGBColor(89, 89, 89)
GOLD = RGBColor(255, 215, 0)

# Complete sensor ablation results
SENSOR_RESULTS = [
    ("frame_r2", 93.50, 77.07, 92.98),
    ("xa_motor", 93.35, 76.59, 93.78),
    ("frame_r1", 93.04, 74.63, 93.74),
    ("y_bed__1", 92.93, 74.80, 93.47),
    ("spindle2", 92.54, 73.01, 92.98),
    ("frame_l2", 92.35, 72.03, 91.32),
    ("frame_l3", 92.28, 71.22, 91.79),
    ("frame_b2", 92.20, 71.22, 93.28),
    ("y_bed__2", 92.12, 71.54, 94.40),
    ("y_bed__3", 91.55, 72.52, 91.90),
    ("y_bed__4", 91.09, 69.92, 93.17),
    ("spindle1", 88.07, 58.54, 91.17),
]

# Complete modality ablation results
MODALITY_ONLY_RESULTS = [
    ("RMS (Computed)", 90.44, 67.48, 40),
    ("Gyroscope", 89.06, 61.30, 72),
    ("Magnetometer", 87.92, 58.21, 72),
    ("Color", 86.69, 55.61, 88),
    ("Accelerometer", 86.50, 51.71, 72),
    ("Environmental", 82.72, 39.67, 72),
]

MODALITY_WITHOUT_RESULTS = [
    ("Without Environmental", 84.36, 45.85, 248),
    ("Without Computed (RMS)", 81.91, 37.40, 280),
    ("Without Accelerometer", 80.96, 37.24, 248),
    ("Without Gyroscope", 81.57, 39.02, 248),
    ("Without Magnetometer", 81.45, 36.42, 248),
    ("Without Color", 79.01, 37.56, 232),
]

# V22 ensemble results
ENSEMBLE_RESULTS = [
    (42, 84.28, 50.57, 85.83),
    (123, 82.98, 50.08, 84.76),
    (456, 83.75, 47.80, 85.37),
]

# Per-modality high correlations
HIGH_CORRELATIONS = [
    ("Pressure", "frame_r1 ↔ frame_b2", "1.000"),
    ("Temperature", "frame_r1 ↔ frame_b2", "1.000"),
    ("Ax", "frame_r1 ↔ frame_b2", "1.000"),
    ("Pressure", "frame_r2 ↔ frame_l2", "1.000"),
    ("RMS", "frame_r2 ↔ y_bed__3", "0.997"),
    ("RMS", "frame_r1 ↔ frame_b2", "0.996"),
    ("My", "frame_r1 ↔ frame_b2", "-0.990"),
    ("Mz", "frame_r1 ↔ frame_b2", "0.988"),
    ("RMS", "frame_r2 ↔ spindle2", "0.986"),
    ("Mx", "frame_r1 ↔ frame_b2", "-0.980"),
    ("Gz", "frame_r1 ↔ frame_b2", "0.973"),
]


def set_cell_text(cell, text, font_size=12, bold=False, color=None, align_center=True):
    """Set text in a table cell with formatting."""
    cell.text = str(text)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        if color:
            paragraph.font.color.rgb = color
        if align_center:
            paragraph.alignment = PP_ALIGN.CENTER


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background gradient effect using shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2), SLIDE_WIDTH, Inches(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.33), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, section_title, section_number=""):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Full blue background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Section number
    if section_number:
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_number
        p.font.size = Pt(80)
        p.font.bold = True
        p.font.color.rgb = RGBColor(100, 150, 200)
        p.alignment = PP_ALIGN.CENTER

    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.33), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullet_points=None):
    """Add a content slide with title and optional bullet points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Bullet points
    if bullet_points:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.33), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            if point.startswith("**") and point.endswith("**"):
                p.text = point.strip("*")
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = DARK_BLUE
                p.space_before = Pt(14)
            elif point.startswith("  - "):
                p.text = point.strip()
                p.font.size = Pt(16)
                p.level = 1
            elif point == "":
                p.text = ""
                p.space_before = Pt(6)
            else:
                p.text = f"• {point}"
                p.font.size = Pt(18)
                p.space_before = Pt(4)

    return slide


def add_table_slide(prs, title, headers, rows, highlight_rows=None, col_widths=None, font_size=11, subtitle=None):
    """Add a slide with a table."""
    slide = add_content_slide(prs, title)

    # Subtitle if provided
    start_top = Inches(1.3)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.4))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = GRAY
        start_top = Inches(1.6)

    # Create table
    num_cols = len(headers)
    num_rows = len(rows) + 1

    left = Inches(0.3)
    width = Inches(12.7)
    row_height = 0.35
    height = Inches(row_height * min(num_rows, 16))

    table = slide.shapes.add_table(num_rows, num_cols, left, start_top, width, height).table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    else:
        col_width = 12.7 / num_cols
        for i in range(num_cols):
            table.columns[i].width = Inches(col_width)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        set_cell_text(cell, header, font_size=font_size + 1, bold=True)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # Data rows
    highlight_rows = highlight_rows or []
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            is_highlight = row_idx in highlight_rows
            set_cell_text(cell, cell_data, font_size=font_size, bold=is_highlight)
            if is_highlight:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 242, 204)
            elif row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)

    return slide


def add_image_slide(prs, title, image_path, caption="", image_width=None):
    """Add a slide with an image."""
    slide = add_content_slide(prs, title)

    if os.path.exists(image_path):
        img_width = Inches(image_width) if image_width else Inches(9)
        img_left = (SLIDE_WIDTH - img_width) / 2
        img_top = Inches(1.4)
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)

    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_two_image_slide(prs, title, image1_path, image2_path, caption1="", caption2=""):
    """Add a slide with two images side by side."""
    slide = add_content_slide(prs, title)

    img_width = Inches(5.8)
    img_top = Inches(1.5)

    if os.path.exists(image1_path):
        slide.shapes.add_picture(image1_path, Inches(0.5), img_top, width=img_width)

    if os.path.exists(image2_path):
        slide.shapes.add_picture(image2_path, Inches(7), img_top, width=img_width)

    if caption1:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(5.8), Inches(0.4))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption1
        p.font.size = Pt(11)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    if caption2:
        cap_box = slide.shapes.add_textbox(Inches(7), Inches(6.5), Inches(5.8), Inches(0.4))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption2
        p.font.size = Pt(11)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_two_column_slide(prs, title, left_content, right_content):
    """Add a slide with two columns."""
    slide = add_content_slide(prs, title)

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.9), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if point.startswith("**"):
            p.text = point.strip("*")
            p.font.bold = True
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_BLUE
            p.space_before = Pt(10)
        elif point == "":
            p.text = ""
        else:
            p.text = f"• {point}"
            p.font.size = Pt(14)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.4), Inches(5.9), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if point.startswith("**"):
            p.text = point.strip("*")
            p.font.bold = True
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_BLUE
            p.space_before = Pt(10)
        elif point == "":
            p.text = ""
        else:
            p.text = f"• {point}"
            p.font.size = Pt(14)

    return slide


def create_presentation():
    """Create the full v22 presentation - V2 Polished."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Base path for figures
    base_path = "/Users/stepheneacuello/Projects/gcode_fingerprinting"

    # ========== SLIDE 1: Title ==========
    add_title_slide(
        prs,
        "G-code Fingerprinting",
        "2K Vocabulary Experiment Results | Version 22 | January 2026"
    )

    # ========== SLIDE 2: Executive Summary ==========
    slide = add_table_slide(
        prs,
        "Executive Summary: Key Results",
        ["Configuration", "Token Acc.", "Seq. Acc.", "Features"],
        [
            ["Best Single Sensor (frame_r2)", "93.50%", "77.1%", "41"],
            ["All 12 Sensors (v22 corrected)", "84.28%", "50.6%", "296"],
            ["All 12 Sensors (without color)", "75.37%", "28.0%", "208"],
            ["Performance Gap (single vs multi)", "-9.22%", "-26.5%", "+255"],
        ],
        highlight_rows=[0, 3],
        font_size=14,
        subtitle="Single sensor outperforms multi-sensor fusion by 9.2%"
    )

    # ========== SECTION: Dataset & Architecture ==========
    add_section_slide(prs, "Dataset & Architecture", "01")

    # ========== SLIDE 3: Dataset Overview ==========
    add_two_column_slide(
        prs,
        "Dataset Overview",
        [
            "**Data Composition**",
            "145 aligned CSV files from real CNC operations",
            "3,160+ training sequences (64 timesteps each)",
            "2,000+ token vocabulary (4-digit precision)",
            "",
            "**Operations (145 files)**",
            "Face operations: 40 files (27.6%)",
            "Pocket operations: 40 files (27.6%)",
            "Adaptive operations: 40 files (27.6%)",
            "Damaged operations: 15 files (10.3%)",
            "150025 variants: 10 files (6.9%)",
        ],
        [
            "**Sensors (12 total)**",
            "Frame: frame_r1, frame_r2, frame_b2, frame_l2, frame_l3",
            "Spindle: spindle1, spindle2",
            "Y-bed: y_bed__1, y_bed__2, y_bed__3, y_bed__4",
            "Motor: xa_motor",
            "",
            "**Modalities per Sensor (17)**",
            "Motion: Ax, Ay, Az, Gx, Gy, Gz, Mx, My, Mz",
            "Environmental: Pressure, Temperature, Proximity",
            "Color: ColorR, ColorG, ColorB, ColorA",
            "Computed: RMS",
            "",
            "**Total Features: 12 × ~25 = 296**",
        ]
    )

    # ========== SLIDE: Preprocessing Pipeline ==========
    add_content_slide(
        prs,
        "Preprocessing Pipeline",
        [
            "**Step 1: Temporal Alignment**",
            "  - Interpolate sensor + motor data onto controller (SR) timestamps",
            "  - Linear interpolation via np.interp() for consistent time base",
            "",
            "**Step 2: Feature Construction (296 features)**",
            "  - 12 sensors x ~17 modalities (IMU, Pressure, Temp, Proximity, Color, RMS)",
            "  - + Motor currents, machine positions, control signals",
            "",
            "**Step 3: G-code Tokenization (2K vocab)**",
            "  - Hybrid tokenizer: literal commands (G0, G1) + bucketed numeric values",
            "  - Canonicalization: strip comments, normalize codes, 4-digit precision",
            "",
            "**Step 4: Windowing**",
            "  - 64 timesteps per window, stride 16 (75% overlap)",
            "  - Each window → sensor tensor [64, 296] + G-code token sequence",
            "",
            "**Step 5: Feature Selection & Cleaning**",
            "  - Remove >50% missing, zero-variance, and |r|>0.95 correlated features",
            "  - Forward-fill NaN, clip outliers (IQR), log-transform skewed features",
        ]
    )

    # ========== SLIDE: Preprocessing Pipeline (cont.) ==========
    add_two_column_slide(
        prs,
        "Preprocessing Pipeline (cont.)",
        [
            "**Step 6: Normalization**",
            "v21: Global StandardScaler on ALL data (leaked!)",
            "v22: StandardScaler on TRAINING data only",
            "Per-feature (each of 296 columns independently)",
            "",
            "**Step 7: Train / Val / Test Split**",
            "70% train: 93 files (2,966 samples)",
            "15% val: 21 files (614 samples)",
            "15% test: 21 files (615 samples)",
            "File-level split (no sample leakage)",
        ],
        [
            "**Step 8: Data Augmentation (training only)**",
            "Sensor noise (Gaussian, sigma=0.01)",
            "Temporal shift (±1 timestep)",
            "Magnitude scaling (0.98-1.02x)",
            "Time warping (non-linear distortion)",
            "Feature dropout (p=0.05 per channel)",
            "Temporal cutout (mask 1-2 timesteps)",
            "Class balancing (3x oversample rare G-codes)",
            "Mixup (blend similar operations, alpha=0.2)",
        ]
    )

    # ========== SLIDE: Experimental Setup Photo ==========
    setup_photo = f"{base_path}/outputs/archive/2024-12_best_runs/sensor_multihead_v3/visualizations/paper/figures/setup.png"
    sensor_setup = f"{base_path}/outputs/archive/2024-12_best_runs/sensor_multihead_v3/visualizations/paper/figures/sensor_setup.png"
    add_two_image_slide(
        prs,
        "Experimental Setup & Sensor Hardware",
        setup_photo,
        sensor_setup,
        "Bantam Explorer Desktop CNC with mounted Arduino Nano 33 BLE Sense sensors",
        "Sensor pipeline: 9-axis IMU, Pressure, Temperature, Proximity, Color (RGBA)"
    )

    # ========== SLIDE: Implementation Details ==========
    add_two_column_slide(
        prs,
        "Implementation Details",
        [
            "**Feature Extraction**",
            "Raw CSV sensor data, 17 modalities/sensor",
            "Columns: {sensor_id}.{modality}",
            "64 timesteps per sequence",
            "",
            "**Normalization**",
            "v21: Global StandardScaler on ALL data (leaked!)",
            "v22: StandardScaler fitted on TRAINING data only",
            "Per-feature normalization (each of 296 features)",
            "Not per-sensor (may reduce natural correlations)",
        ],
        [
            "**Train / Val / Test Split**",
            "70% train: 93 files (2,966 samples)",
            "15% validation: 21 files (614 samples)",
            "15% test: 21 files (615 samples)",
            "File-level split (no sample leakage)",
            "",
            "**Correlation Computation**",
            "v21: Averaged all 17 modalities per sensor first",
            "v22: Per-modality correlation computed separately",
            "Per-modality approach reveals true physical correlations",
        ]
    )

    # ========== SLIDE: Architecture Diagram ==========
    arch_path = f"{base_path}/outputs/archive/2024-12_best_runs/sensor_multihead_v3/visualizations/paper/figures/simplified_architecture.png"
    add_image_slide(
        prs,
        "Two-Stage Architecture",
        arch_path,
        "Stage 1: MM-DTAE-LSTM Encoder (100% operation classification) → Stage 2: Transformer Multi-Head Decoder",
        image_width=11
    )

    # ========== SLIDE: Training Curves ==========
    train_curves = f"{base_path}/outputs/archive/2024-12_best_runs/sensor_multihead_v3/visualizations/paper/figures/training_curves_4panel.png"
    add_image_slide(
        prs,
        "Training Dynamics: Encoder & Decoder",
        train_curves,
        "Encoder converges to 100% accuracy by epoch 30. Decoder reaches ~90% with curriculum learning phases P1/P2/P3.",
        image_width=7
    )

    # ========== SLIDE: Encoder Performance ==========
    encoder_cm = f"{base_path}/outputs/production/encoder/confusion_matrix.png"
    add_image_slide(
        prs,
        "Stage 1: Encoder Operation Classification (100% Accuracy)",
        encoder_cm,
        "The encoder perfectly classifies all 9 operation types from sensor data",
        image_width=8
    )

    # ========== SECTION: V22 Improvements ==========
    add_section_slide(prs, "V22 Improvements", "02")

    # ========== SLIDE 6: Data Leakage Fix ==========
    slide = add_table_slide(
        prs,
        "V22 Improvement #1: Data Leakage Fix",
        ["Metric", "v21 (With Leakage)", "v22 (Corrected)", "Improvement"],
        [
            ["Multi-sensor token acc", "80.69%", "84.28%", "+3.59%"],
            ["Multi-sensor seq acc", "35.4%", "50.6%", "+15.2%"],
            ["Single sensor token acc", "93.50%", "93.50%", "unchanged"],
            ["Performance gap", "12.81%", "9.22%", "narrowed"],
        ],
        highlight_rows=[0, 1],
        font_size=14,
        subtitle="Issue: StandardScaler fitted on all data → Fix: Fitted ONLY on training data"
    )

    # ========== SLIDE 7: Model Stability ==========
    ensemble_rows = [[str(s), f"{t:.2f}%", f"{sq:.2f}%", f"{v:.2f}%"]
                     for s, t, sq, v in ENSEMBLE_RESULTS]
    ensemble_rows.append(["Mean", "83.67%", "49.48%", "85.32%"])
    ensemble_rows.append(["Std Dev", "±0.66%", "±1.52%", "±0.54%"])

    add_table_slide(
        prs,
        "V22 Improvement #2: Model Stability Across Seeds",
        ["Seed", "Token Acc.", "Seq. Acc.", "Val. Token"],
        ensemble_rows,
        highlight_rows=[3, 4],
        font_size=14,
        subtitle="Very low variance (±0.66%) demonstrates robust, well-tuned architecture"
    )

    # ========== SLIDE 8: Color Ablation ==========
    add_table_slide(
        prs,
        "V22 Improvement #3: Color Sensors Are Essential",
        ["Configuration", "Features", "Token Acc.", "Seq. Acc."],
        [
            ["With Color (RGBA)", "296", "84.28%", "50.57%"],
            ["Without Color", "208", "75.37%", "27.97%"],
            ["Difference", "-88", "-8.91%", "-22.60%"],
        ],
        highlight_rows=[2],
        font_size=14,
        subtitle="Despite R,G,B,A being 97%+ correlated internally, they provide essential cross-modal information"
    )

    # ========== SECTION: Sensor Analysis ==========
    add_section_slide(prs, "Sensor Analysis", "03")

    # ========== SLIDE 9: Complete Sensor Rankings ==========
    sensor_rows = [[str(i+1), name, f"{tok:.2f}%", f"{seq:.2f}%", f"{val:.2f}%"]
                   for i, (name, tok, seq, val) in enumerate(SENSOR_RESULTS)]
    sensor_rows.append(["-", "All 12 Sensors (v22)", "84.28%", "50.6%", "85.83%"])

    add_table_slide(
        prs,
        "Complete Single Sensor Performance (All 12 Sensors)",
        ["Rank", "Sensor", "Test Token", "Test Seq.", "Val Token"],
        sensor_rows,
        highlight_rows=[0, 12],
        font_size=10,
        subtitle="Every single sensor (88-93%) outperforms all sensors combined (84%)"
    )

    # ========== SLIDE 10: Sensor Correlation Heatmap ==========
    sensor_corr = f"{base_path}/outputs/analysis/sensor_correlation_heatmap.png"
    add_image_slide(
        prs,
        "Sensor Correlation Matrix (Averaged Across Modalities)",
        sensor_corr,
        "Original analysis: Max |r| = 0.56, suggesting sensors are NOT redundant. But this masked per-modality patterns...",
        image_width=8
    )

    # ========== SLIDE 11: Per-Modality Correlation ==========
    per_mod_corr = f"{base_path}/outputs/jan23_followup/correlations/per_modality_correlation_summary.png"
    add_image_slide(
        prs,
        "Per-Modality Sensor Correlations (Professor Sodhi's Insight)",
        per_mod_corr,
        "Per-modality analysis reveals frame sensors ARE highly correlated (r=1.00) for Pressure, Temperature, Ax",
        image_width=11
    )

    # ========== SLIDE 12: High Correlation Table ==========
    corr_rows = [[mod, pair, r] for mod, pair, r in HIGH_CORRELATIONS]
    add_table_slide(
        prs,
        "High Per-Modality Correlations (|r| > 0.97)",
        ["Modality", "Sensor Pair", "Correlation"],
        corr_rows,
        highlight_rows=[0, 1, 2, 3],
        font_size=12,
        subtitle="Environmental and RMS modalities show near-perfect correlation across frame sensors"
    )

    # ========== SLIDE: Sensor Location Importance ==========
    loc_bars = f"{base_path}/outputs/archive/2024-12_best_runs/sensor_multihead_v3/visualizations/paper/figures/location_importance_bars.png"
    add_image_slide(
        prs,
        "Sensor Location Importance (Leave-One-Out)",
        loc_bars,
        "Accuracy drop when each sensor location is removed. y_bed_3 and frame_r1 are most important for multi-sensor model.",
        image_width=8
    )

    # ========== SLIDE: Modality × Location Heatmap ==========
    loc_heatmap = f"{base_path}/outputs/archive/2024-12_best_runs/sensor_multihead_v3/visualizations/paper/figures/sensor_location_heatmap.png"
    add_image_slide(
        prs,
        "Modality x Location Interaction Matrix",
        loc_heatmap,
        "Most critical: Proximity at frame_r1 (-2.77% accuracy drop). Color and Motor Current contribute least.",
        image_width=10
    )

    # ========== SECTION: Modality Analysis ==========
    add_section_slide(prs, "Modality Analysis", "04")

    # ========== SLIDE 13: Single Modality Performance ==========
    mod_rows = [[name, f"{tok:.2f}%", f"{seq:.2f}%", str(feat)]
                for name, tok, seq, feat in MODALITY_ONLY_RESULTS]
    mod_rows.append(["All Modalities (v22)", "84.28%", "50.6%", "296"])

    add_table_slide(
        prs,
        "Single Modality Performance",
        ["Modality Only", "Token Acc.", "Seq. Acc.", "Features"],
        mod_rows,
        highlight_rows=[0, 6],
        font_size=13,
        subtitle="RMS alone (90.44%) outperforms all modalities combined (84.28%)!"
    )

    # ========== SLIDE 14: Leave-One-Out Modality ==========
    without_rows = [[name, f"{tok:.2f}%", f"{seq:.2f}%", str(feat)]
                    for name, tok, seq, feat in MODALITY_WITHOUT_RESULTS]

    add_table_slide(
        prs,
        "Leave-One-Out Modality Ablation",
        ["Configuration", "Token Acc.", "Seq. Acc.", "Features"],
        without_rows,
        highlight_rows=[0, 5],
        font_size=13,
        subtitle="Removing environmental helps (+3.7%), removing color hurts most (-5.3%)"
    )

    # ========== SLIDE 15: Modality Correlation Heatmap ==========
    mod_corr = f"{base_path}/outputs/analysis/modality_correlation_heatmap.png"
    add_image_slide(
        prs,
        "Modality Correlation Matrix",
        mod_corr,
        "Color channels R,G,B,A are 97%+ correlated. Ax↔Gz, Gz↔Pressure also highly correlated.",
        image_width=9
    )

    # ========== SECTION: Confusion Matrices ==========
    add_section_slide(prs, "Model Performance Analysis", "05")

    # ========== SLIDE 16: Type Confusion Matrix ==========
    type_cm = f"{base_path}/outputs/real_confusion_matrices/confusion_type.png"
    add_image_slide(
        prs,
        "Token Type Classification (4 Classes)",
        type_cm,
        "Numeric: 99%, Special: 100%, Parameter: 89%, Command: 80%",
        image_width=7
    )

    # ========== SLIDE 17: Command Confusion Matrix ==========
    cmd_cm = f"{base_path}/outputs/real_confusion_matrices/confusion_command.png"
    add_image_slide(
        prs,
        "G-code Command Classification",
        cmd_cm,
        "G0: 89%, G1: 80%, G53: 45% (rare command confusion)",
        image_width=7
    )

    # ========== SLIDE 18: Parameter Confusion Matrix ==========
    param_cm = f"{base_path}/outputs/real_confusion_matrices/confusion_parameter.png"
    add_image_slide(
        prs,
        "Parameter Type Classification",
        param_cm,
        "F: 100%, R: 100%, X: 94%, Z: 99%, Y: 71% (X/Y confusion in similar motions)",
        image_width=7
    )

    # ========== SLIDE: Digit Accuracy Heatmap ==========
    digit_heatmap = f"{base_path}/outputs/archive/2024-12_best_runs/sensor_multihead_v3/visualizations/paper/figures/digit_accuracy_heatmap.png"
    add_image_slide(
        prs,
        "Per-Digit Prediction Accuracy",
        digit_heatmap,
        "Accuracy degrades with digit precision: Sign (98-99%) → D1/10s (94-96%) → D6/.0001 (79-82%). Z-axis easiest to predict.",
        image_width=7
    )

    # ========== SECTION: Architecture Ablations & Baselines ==========
    add_section_slide(prs, "Architecture Ablations & Baselines", "06")

    # ========== SLIDE: Architecture Ablations Table ==========
    add_table_slide(
        prs,
        "Architecture & Hyperparameter Ablations",
        ["Configuration", "Token Acc.", "Change vs Medium"],
        [
            ["Medium model (baseline)", "78.62%", "-"],
            ["XLarge model", "78.47%", "-0.15%"],
            ["Augmentation prob 0.5", "78.32%", "-0.30%"],
            ["Curriculum learning ON", "77.82%", "-0.80%"],
            ["No augmentation", "77.71%", "-0.91%"],
            ["No operation conditioning", "76.29%", "-2.33%"],
            ["Base model", "76.25%", "-2.37%"],
            ["No positional encoding", "76.18%", "-2.44%"],
            ["Tiny model", "72.97%", "-5.65%"],
            ["LSTM only (no transformer)", "70.82%", "-7.80%"],
        ],
        highlight_rows=[0, 9],
        font_size=12,
        subtitle="Model capacity (tiny → xlarge) only changes accuracy by ~5%. Bottleneck is NOT model size."
    )

    # ========== SLIDE: Baselines ==========
    add_table_slide(
        prs,
        "Baseline Comparisons",
        ["Model", "Token Acc.", "Seq. Acc.", "Notes"],
        [
            ["Random guess", "0.5%", "0.0%", "1/2000 tokens"],
            ["Majority class", "23.74%", "0.0%", "Always predict most common"],
            ["LSTM only (no transformer)", "70.82%", "20.5%", "Encoder only"],
            ["Full model (all sensors, v22)", "84.28%", "50.6%", "296 features"],
            ["Single best sensor (frame_r2)", "93.50%", "77.1%", "41 features"],
        ],
        highlight_rows=[4],
        font_size=13,
        subtitle="Clear progression from trivial baselines to full model. Single sensor still best."
    )

    # ========== SLIDE: Statistical Comparison ==========
    stat_comp = f"{base_path}/outputs/archive/2024-12_best_runs/sensor_multihead_v3/visualizations/paper/figures/statistical_comparison.png"
    add_image_slide(
        prs,
        "Statistical Comparison: Our Model vs Baselines",
        stat_comp,
        "Cross-validation score distributions (left) and Cohen's d effect sizes (right). All differences p < 0.001 (***). Our model (90%+) significantly outperforms all baselines.",
        image_width=11
    )

    # ========== SECTION: Diagnosis & Conclusions ==========
    add_section_slide(prs, "Diagnosis & Conclusions", "07")

    # ========== SLIDE 19: Why Multi-Sensor Underperforms ==========
    add_content_slide(
        prs,
        "Why Does Multi-Sensor Fusion Underperform?",
        [
            "**1. Sample Complexity (Curse of Dimensionality)**",
            "  - With 41 features: 3,160 samples may be sufficient",
            "  - With 296 features: samples are sparse in feature space",
            "",
            "**2. High Redundancy (Per-Modality Analysis)**",
            "  - Pressure, Temperature, Ax: r=1.00 across frame sensors",
            "  - RMS: r=0.99 across multiple sensor pairs",
            "  - Model receives same information multiple times",
            "",
            "**3. Feature Interference**",
            "  - Correlated features → conflicting gradients",
            "  - Attention mechanism diluted across similar inputs",
            "",
            "**4. The Color Paradox**",
            "  - R,G,B,A 97%+ correlated internally (redundant)",
            "  - Yet removing color drops accuracy by 8.9%",
            "  - Color provides CROSS-MODAL information not in other sensors",
        ]
    )

    # ========== SLIDE 20: V22 Summary Table ==========
    add_table_slide(
        prs,
        "Summary: V21 vs V22 Results",
        ["Finding", "V21", "V22", "Change"],
        [
            ["Multi-sensor token acc", "80.69%", "84.28%", "+3.59%"],
            ["Multi-sensor seq acc", "35.4%", "50.6%", "+15.2%"],
            ["Single sensor token acc", "93.50%", "93.50%", "unchanged"],
            ["Performance gap", "12.81%", "9.22%", "narrowed"],
            ["Model variance (3 seeds)", "N/A", "±0.66%", "stable"],
            ["Color contribution", "N/A", "+8.91%", "essential"],
            ["Per-modality correlation", "masked", "r=1.00", "confirmed"],
        ],
        highlight_rows=[0, 1, 4, 5, 6],
        font_size=13
    )

    # ========== SLIDE 21: Key Conclusions ==========
    add_content_slide(
        prs,
        "Key Conclusions",
        [
            "**1. Single Sensor Achieves Best Performance (93.50%)**",
            "  - frame_r2, xa_motor, frame_r1 all viable choices",
            "  - 9.2% better than multi-sensor fusion",
            "",
            "**2. V22 Methodology Fixes Improved Multi-Sensor**",
            "  - Token accuracy: 80.69% → 84.28% (+3.59%)",
            "  - Sequence accuracy: 35.4% → 50.6% (+15.2%)",
            "",
            "**3. Color Sensors Essential Despite Internal Correlation**",
            "  - Removing RGBA drops accuracy by 8.91%",
            "  - Provides unique cross-modal information",
            "",
            "**4. Professor Sodhi's Hypothesis Confirmed**",
            "  - Frame sensors highly correlated per-modality (r=1.00)",
            "  - Original averaging analysis masked true redundancy",
            "",
            "**5. Model is Robust and Well-Tuned**",
            "  - ±0.66% variance across 3 random seeds",
        ]
    )

    # ========== SLIDE 22: Recommended Path Forward ==========
    add_two_column_slide(
        prs,
        "Recommended Path Forward",
        [
            "**Option A: Single Best Sensor**",
            "Accuracy: 93.50% | Features: 41",
            "Pros: Best performance, simplest, fastest",
            "Cons: May not generalize across machines",
            "",
            "**Option B: Top 2-3 Non-Redundant Sensors**",
            "Expected: 90-92% | Features: 82-123",
            "Pros: More robust, complementary info",
            "Cons: Need to identify non-redundant pairs",
        ],
        [
            "**Option C: All Sensors + Modality Selection**",
            "Expected: 86-88% | Features: ~150",
            "Pros: Uses all spatial positions",
            "Cons: Still more complex than needed",
            "",
            "**Recommended Next Steps**",
            "1. Test sensor combinations (e.g., frame_r2 + xa_motor)",
            "2. Remove redundant modalities (keep one of P/T/Ax)",
            "3. Try per-sensor encoders before fusion",
            "4. Collect more data if multi-sensor required",
        ]
    )

    # ========== SLIDE 23: Experiment Statistics ==========
    add_two_column_slide(
        prs,
        "Experiment Statistics",
        [
            "**Total Experiments**",
            "290+ W&B runs",
            "71+ ablation studies",
            "~50 GPU hours training",
            "",
            "**Ablation Categories**",
            "12 single-sensor experiments",
            "12 leave-one-sensor-out",
            "6 single-modality experiments",
            "6 leave-one-modality-out",
            "35+ architecture ablations",
            "3 ensemble seed runs",
            "1 no-color ablation",
        ],
        [
            "**Key Checkpoints**",
            "Best single: sensor_ablations/only_frame_r2/",
            "Best multi (v22): jan23_followup/no_leakage/",
            "Ensemble: jan23_followup/ensemble/",
            "No-color: jan23_followup/no_color/",
            "",
            "**Reproducibility**",
            "Seeds: 42, 123, 456",
            "Config: configs/best_lambda_2k_vocab.json",
            "All logged to Weights & Biases",
        ]
    )

    # ========== SLIDE: Follow-Up Item Status ==========
    add_two_column_slide(
        prs,
        "Email Follow-Up Item Status",
        [
            "**Completed**",
            "[DONE] Data leakage fix (train-only scaler)",
            "[DONE] Model stability (3 seeds: ±0.66%)",
            "[DONE] Color channel ablation (+8.91%)",
            "[DONE] Per-modality correlation analysis",
            "[DONE] Sensor correlation matrix",
            "[DONE] Single-sensor rankings (all 12)",
            "[DONE] Single-modality rankings (all 6)",
            "[DONE] Architecture ablations (10 configs)",
            "[DONE] Statistical comparison (Cohen's d)",
        ],
        [
            "**Not Yet Started**",
            "[ ] Top 2-3 sensor combinations",
            "[ ] PCA / feature selection on 296 features",
            "[ ] Embedding-level sensor fusion",
            "[ ] Per-sensor normalization comparison",
            "[ ] Top-k accuracy analysis",
            "",
            "**⚠ Architecture Update Pending**",
            "Diagrams show Bi-LSTM + Attn Pool;",
            "code uses unidirectional LSTM + Mean Pool.",
            "Once approved, ALL experiments must be re-run.",
        ]
    )

    # ========== SLIDE: Thank You ==========
    add_title_slide(
        prs,
        "Questions?",
        "Stephen Eacuello | seacuello@uri.edu | January 2026"
    )

    # Save presentation
    output_path = f"{base_path}/outputs/presentations/gcode_fingerprinting_v22_results_v2.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    return output_path


if __name__ == "__main__":
    create_presentation()
